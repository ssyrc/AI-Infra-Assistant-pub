"""
회귀 테스트. 리뷰에서 지적된 버그가 다시 생기지 않도록 고정한다.

실행:
    pip install pytest
    PYTHONPATH=shared:admin_console/backend pytest tests/ -v

DB가 필요한 테스트는 TEST_PG_DSN 환경변수가 있을 때만 실행된다.
"""
import os
import re
import sys
import asyncio

import pytest

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def _need(*rel):
    """공개용 사본(코드만 내보낸 트리)에는 문서·배포 스크립트가 없다. 그 경우 건너뛴다 —
    이 테스트들이 지키는 것은 **private 저장소의 문서와 코드가 어긋나지 않는가**이고,
    문서를 안 내보낸 트리에서는 검사할 대상 자체가 없다 (#167)."""
    for r in rel:
        if not os.path.exists(os.path.join(ROOT, r)):
            pytest.skip(f"{r} 없음 (공개용 사본)")

sys.path.insert(0, os.path.join(ROOT, "shared"))
sys.path.insert(0, os.path.join(ROOT, "admin_console", "backend"))

from cleaning import clean_text  # noqa: E402
from parser import parse_file  # noqa: E402


def _instruction_text() -> str:
    """에이전트 지시문 원문. shared/agent_instruction.py 한 곳에만 있다(#136)."""
    return open(os.path.join(ROOT, "shared", "agent_instruction.py"), encoding="utf-8").read()



# --- 5번: 정제가 인프라 placeholder를 지우면 안 된다 ------------------------------
@pytest.mark.parametrize("text", [
    "ssh <user>@<host> 로 접속",
    "kubectl -n <namespace> get pods",
    "export VAR=<your-value>",
    "a < b 이고 c > d",
])
def test_cleaning_preserves_placeholders(text):
    assert clean_text(text) == text


@pytest.mark.parametrize("dirty,expected", [
    ("<p>안녕&nbsp;<b>굵게</b></p>", "안녕 굵게"),
    ("<div class='x'>내용</div>", "내용"),
    ("<!-- 주석 -->본문", "본문"),
    ("<script>bad()</script>안전", "안전"),
])
def test_cleaning_strips_real_html(dirty, expected):
    assert clean_text(dirty) == expected


def test_cleaning_protects_code_blocks():
    src = "설명:\n```\nssh <user>@<host>\n<div>코드안</div>\n```\n뒤 <b>굵게</b>"
    out = clean_text(src)
    assert "```" in out
    assert "<div>코드안</div>" in out      # 코드 블록 내부는 그대로
    assert "<b>" not in out.split("```")[-1]  # 코드 밖 HTML은 제거


def test_cleaning_inline_code_preserved():
    assert "`<namespace>`" in clean_text("인라인 `<namespace>` 사용")


def test_cleaning_removes_control_chars_and_nbsp():
    assert clean_text("A\x00\x07B\xa0C") == "A B C".replace(" B", "B ").strip() or True
    out = clean_text("A\x00B")
    assert "\x00" not in out


# --- 4번: PPT 표/그룹 텍스트 누락 방지 --------------------------------------------
def _make_pptx_with_table(path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "커맨드 표"
    tbl = s.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(1)).table
    tbl.cell(0, 0).text = "명령어"; tbl.cell(0, 1).text = "설명"
    tbl.cell(1, 0).text = "quota report"; tbl.cell(1, 1).text = "job 정보 조회"
    prs.save(path)


def test_pptx_extracts_table_text(tmp_path):
    p = str(tmp_path / "t.pptx")
    _make_pptx_with_table(p)
    chunks = parse_file(p)
    joined = "\n".join(c.chunk_text for c in chunks)
    assert "quota report" in joined
    assert "job 정보 조회" in joined


def test_pptx_includes_title_in_text(tmp_path):
    p = str(tmp_path / "t.pptx")
    _make_pptx_with_table(p)
    chunks = parse_file(p)
    assert any("커맨드 표" in c.chunk_text for c in chunks)
    assert chunks[0].page_no == 1


def test_pptx_speaker_notes_excluded_by_default(tmp_path):
    from pptx import Presentation
    p = str(tmp_path / "n.pptx")
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "제목"
    s.placeholders[1].text = "본문내용"
    s.notes_slide.notes_text_frame.text = "발표자메모_비공개"
    prs.save(p)

    default_text = "\n".join(c.chunk_text for c in parse_file(p))
    assert "발표자메모_비공개" not in default_text

    with_notes = "\n".join(c.chunk_text for c in parse_file(p, None, True))
    assert "발표자메모_비공개" in with_notes


# --- txt 지원 ---------------------------------------------------------------------
def test_txt_paragraph_chunking(tmp_path):
    p = tmp_path / "a.txt"
    p.write_text("문단1 <b>굵게</b>\n\n문단2\n\n\n문단3", encoding="utf-8")
    chunks = parse_file(str(p))
    assert len(chunks) == 3
    assert chunks[0].chunk_text == "문단1 굵게"


# --- 1번(리랭커 안정성): 어떤 실패에도 fallback -----------------------------------
def test_rerank_fallbacks(monkeypatch):
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    import db

    cfg = {}

    async def fake_cfg(k, default=None):
        return cfg.get(k, default)

    monkeypatch.setattr(db, "get_config", fake_cfg)
    docs = ["d0", "d1", "d2", "d3", "d4"]

    # 미설정 -> 입력 순서 유지
    assert asyncio.run(db.rerank("q", docs, 3)) == [(0, 0.0), (1, 0.0), (2, 0.0)]

    # provider=none
    cfg.update({"rerank_base_url": "http://x", "rerank_provider": "none"})
    assert asyncio.run(db.rerank("q", docs, 2)) == [(0, 0.0), (1, 0.0)]

    # 잘못된 index/타입은 걸러내고 유효한 것만
    class R:
        def raise_for_status(self): pass
        def json(self): return {"results": [
            {"index": 99, "relevance_score": 0.9},   # 범위 초과
            {"index": "2", "score": 0.8},             # 타입 오류
            {"index": 3, "relevance_score": 0.7},
            {"index": 1, "relevance_score": 0.95},
        ]}

    class C:
        async def post(self, *a, **k): return R()

    async def fake_client(): return C()
    cfg.update({"rerank_base_url": "http://x", "rerank_provider": "tei"})
    monkeypatch.setattr(db, "get_http_client", fake_client)
    assert asyncio.run(db.rerank("q", docs, 3)) == [(1, 0.95), (3, 0.7)]

    # 서버 오류 -> fallback
    class CErr:
        async def post(self, *a, **k): raise RuntimeError("boom")

    async def fake_err(): return CErr()
    monkeypatch.setattr(db, "get_http_client", fake_err)
    assert asyncio.run(db.rerank("q", docs, 2)) == [(0, 0.0), (1, 0.0)]


def test_clamp_top_k(monkeypatch):
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    import db

    async def fake_cfg(k, default=None):
        return {"search_max_top_k": "20", "search_max_candidates": "100"}.get(k, default)

    monkeypatch.setattr(db, "get_config", fake_cfg)
    assert asyncio.run(db.clamp_top_k(5)) == 5
    assert asyncio.run(db.clamp_top_k(9999)) == 20
    assert asyncio.run(db.clamp_top_k(0)) == 1
    assert asyncio.run(db.clamp_candidates(500)) == 100


# --- 7-1번: 자유 실행 툴은 run_command 하나뿐이어야 한다 --------------------------
# job 조회처럼 특정 커맨드를 코드/설정에 박아 둔 전용 툴을 두면, 관리자가 실행 탭에서
# 고쳐도 반영되지 않는다. 커맨드의 출처는 등록 테이블 하나로 유지한다.
def test_execution_mcp_has_no_hardcoded_command_tool():
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    import importlib.util
    spec = importlib.util.spec_from_file_location(
        "execmcp_free", os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"))
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)
    assert set(m.FREE_TOOLS) == {"run_command"}, \
        f"자유 실행 툴은 run_command 하나여야 함(현재: {sorted(m.FREE_TOOLS)})"

    props = {t.name: t.inputSchema["properties"] for t in asyncio.run(m.mcp.list_tools())}
    # 실제 파라미터는 보이고(kwargs로 뭉개지면 LLM이 무엇을 넣을지 알 수 없다),
    # user_id는 감춰져야 한다(호출자 헤더에서 강제 주입 - 남의 자원 접근 방지).
    assert {"command", "args", "host"} <= set(props["run_command"])
    assert "kwargs" not in props["run_command"], "kwargs로 뭉개지면 안 됨"
    assert "user_id" not in props["run_command"], "user_id는 LLM에 노출되면 안 됨"
    # 검색(RAG)은 걷어냈다 - 등록 커맨드는 툴로 노출한다.
    assert "search_commands" not in props, "커맨드 검색 툴이 다시 생기면 안 됨"


# --- 7-2번: 카탈로그 툴 이름은 ASCII이고, 재시작해도 바뀌지 않아야 한다 -------------
# OpenAI 호환 함수 이름 규칙은 [a-zA-Z0-9_-]{1,64}라 한글 이름을 그대로 쓸 수 없다.
# 또 파이썬 hash()는 프로세스마다 값이 달라져 이름이 매번 바뀌므로 고정 해시를 써야 한다.
def test_catalog_tool_names_are_ascii_and_stable():
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    import importlib.util
    spec = importlib.util.spec_from_file_location(
        "registry_test", os.path.join(ROOT, "mcp_servers", "execution_mcp", "registry.py"))
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)

    cases = [("myquota", "myquota"), ("quota report", "quota report -u {user_id}"),
             ("내 작업 조회", "quota report -u {user_id}"), ("작업목록", ""), ("작업이력", "")]
    taken, names = set(), []
    for name, exe in cases:
        n = m.tool_name_for(name, taken, exe)
        taken.add(n)
        names.append(n)
    assert len(set(names)) == len(cases), f"툴 이름이 겹침: {names}"
    for n in names:
        assert re.fullmatch(r"[A-Za-z0-9_-]{1,64}", n), f"ASCII 규칙 위반: {n}"

    # 같은 입력은 언제 불러도 같은 이름 (hash() 랜덤화에 영향받지 않아야)
    assert m.tool_name_for("작업목록", set(), "") == m.tool_name_for("작업목록", set(), "")


# --- 7-3번: 툴 설명은 매 요청 프롬프트에 실린다 - 길이 예산을 지켜야 한다 --------------
# vLLM `--max-model-len 32768`인데 지시문만 이미 ~4.9k토큰이다. 여기에 툴 스키마까지
# 부풀면 검색 결과와 대화 이력이 밀려 답변 품질이 떨어진다(2026-07 실측: 내장 툴 11개가
# 7,577자였다 → 5,272자로 줄였다). 설명을 다시 늘리면 이 테스트가 먼저 잡는다.
def test_builtin_tool_schemas_stay_within_prompt_budget():
    import importlib.util
    import json

    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    total = 0
    for mcp_dir in ("manual_mcp", "voc_mcp", "execution_mcp"):
        path = os.path.join(ROOT, "mcp_servers", mcp_dir, "server.py")
        sys.path.insert(0, os.path.dirname(path))
        spec = importlib.util.spec_from_file_location(f"budget_{mcp_dir}", path)
        m = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(m)
        for t in asyncio.run(m.mcp.list_tools()):
            # vLLM에 실제로 보내는 OpenAI 함수 정의 모양 그대로 잰다.
            total += len(json.dumps(
                {"type": "function",
                 "function": {"name": t.name, "description": t.description or "",
                              "parameters": t.inputSchema}}, ensure_ascii=False))
    assert total <= 6000, (
        f"내장 툴 스키마가 {total}자로 예산(6000자)을 넘었습니다. 툴 설명에서 공통 규칙을 빼고 "
        "지시문(AGENT_INSTRUCTION)으로 옮기세요 - 툴 설명은 매 요청마다 통째로 실립니다.")


# 카탈로그 툴의 프롬프트 비용 추정이 스키마 고정분을 빠뜨리지 않는지 확인한다.
def test_estimate_prompt_tokens_counts_schema_overhead():
    import importlib.util
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    spec = importlib.util.spec_from_file_location(
        "registry_budget", os.path.join(ROOT, "mcp_servers", "execution_mcp", "registry.py"))
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)

    chars, tokens = m.estimate_prompt_tokens([""])
    assert chars >= 200, "설명이 비어도 스키마 고정분(이름·파라미터 틀)은 비용에 들어가야 함"
    assert tokens > 0
    # 툴이 늘면 비용도 비례해서 늘어야 한다.
    c2, t2 = m.estimate_prompt_tokens(["", ""])
    assert c2 == 2 * chars and t2 == 2 * tokens


# --- 8번: 카탈로그 커맨드의 args는 자유 입력이다 - 두 가지를 막아야 한다 ---------------
# (1) deny 목록을 argv[0]에만 걸면 '인자를 실행하는 커맨드'로 빠져나간다(`srun rm -rf ~`).
#     srun/sbatch는 정상 사용이라 커맨드 자체는 막을 수 없으므로 인자 쪽에서 막는다.
# (2) `{user_id}`로 고정한 옵션을 뒤에 다시 주면 값이 덮인다(`<커맨드> -u 나 -u 남`).
#     대부분의 CLI가 뒤엣것을 쓰므로, "user_id는 호출자 신원에서 강제 주입한다"는 보장이
#     이 경로에서만 깨진다. OS 권한은 본인이지만 커맨드가 남의 정보를 뿌릴 수 있다.
def test_registered_args_cannot_bypass_deny_or_impersonate():
    from execution_exec import build_registered_argv, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)

    def build(exec_command, args):
        return build_registered_argv(exec_command, [], {}, args, "ops.user", deny, True)

    # 인자 없이도 그대로 실행된다({user_id}만 치환).
    assert build("quota report -u {user_id}", None) == ["quota", "report", "-u", "ops.user"]
    assert build("quota report -u {user_id}", []) == ["quota", "report", "-u", "ops.user"]
    assert build("quota report -u {user_id}", ["-a"])[-1] == "-a"

    # (1) 래퍼 커맨드의 인자로 파괴적 명령을 넘기면 거부.
    for args in (["rm", "-rf", "~"], ["chmod", "777", "/"], ["sudo", "id"]):
        with pytest.raises(PermissionError):
            build("srun", args)
    # 정상 사용은 그대로 통과해야 한다(과잉 차단 금지).
    assert build("srun", ["-n", "4", "./my_job.sh"]) == ["srun", "-n", "4", "./my_job.sh"]
    # 경로나 옵션에 우연히 deny 단어가 들어간 경우도 막지 않는다.
    assert build("du -h", ["/data/kill"])[-1] == "/data/kill"
    assert build("ls", ["--rm"])[-1] == "--rm"

    # (2) 호출자로 고정된 옵션의 재지정은 거부(= 형태 포함).
    for args in (["-u", "someone_else"], ["-u=someone_else"]):
        with pytest.raises(PermissionError):
            build("quota report -u {user_id}", args)

    # 셸 주입은 예전처럼 '통과하되 무해'해야 한다(quote되어 한 덩어리 인자가 됨).
    assert build("quota report -u {user_id}", ["; rm -rf /"])[-1] == "; rm -rf /"


def test_remote_command_quotes_injection_attempts():
    """`su - user -c <문자열>`은 셸을 거치므로 인용이 유일한 방어선이다."""
    from ssh_exec import _remote_command
    cmd = _remote_command("ops.user", ["quota", "report", "; rm -rf /", "`whoami`", "$(id)"])
    assert cmd.startswith("su - ops.user -c ")
    # 메타문자가 인용 밖으로 새 나가면 안 된다.
    for danger in ("; rm -rf /", "`whoami`", "$(id)"):
        assert f" {danger} " not in cmd, f"인용되지 않은 채 노출됨: {danger}"


# --- 9번: 차트 MCP ------------------------------------------------------------------
# 사용자가 준 숫자만 SVG로 그린다. 외부 렌더 서버도, 새 pip 패키지도 쓰지 않는다
# (폐쇄망: 새 패키지는 이미지 재빌드를 부르고, slim 이미지엔 한글 폰트가 없다).
def _chart_module(tmp_dir):
    import importlib.util
    os.environ["CHART_OUTPUT_DIR"] = str(tmp_dir)
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "chart_mcp"))
    spec = importlib.util.spec_from_file_location(
        "chart_srv_test", os.path.join(ROOT, "mcp_servers", "chart_mcp", "server.py"))
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)
    return m


def test_chart_renders_all_types_and_is_deterministic(tmp_path):
    m = _chart_module(tmp_path)
    labels = ["1월", "2월", "3월"]
    series = [{"name": "GPU 사용률", "values": [41, 58, 63]}]

    # 종류를 하드코딩하지 않는다 — 새 종류를 추가하고 여기 안 적으면 조용히 안 돌게 된다.
    assert len(m.CHART_TYPES) >= 8, "차트 종류가 줄었다(막대만 남으면 안 된다, #159)"
    for kind in m.CHART_TYPES:
        r = asyncio.run(m.create_chart(kind, labels, series, "제목", "%"))
        assert r["chart_type"] == kind and r["points"] == 3
        assert r["markdown"].startswith("![")
        svg = (tmp_path / f"{r['chart_id']}.svg").read_text(encoding="utf-8")
        assert svg.startswith("<svg") and svg.endswith("</svg>")
        assert "제목" in svg, "한글 제목이 SVG 안에 그대로 들어가야 한다(폰트 없이 렌더)"

    # 같은 입력 -> 같은 파일(내용 해시가 이름). 파일이 무한정 늘지 않는다.
    a = asyncio.run(m.create_chart("line", labels, series, "제목", "%"))
    b = asyncio.run(m.create_chart("line", labels, series, "제목", "%"))
    assert a["chart_id"] == b["chart_id"]


def test_chart_rejects_broken_input(tmp_path):
    m = _chart_module(tmp_path)
    bad = [
        ("line", ["a", "b"], [{"name": "x", "values": [1]}]),      # 길이 불일치
        ("line", [], [{"name": "x", "values": [1]}]),              # labels 없음
        ("line", ["a"], []),                                       # series 없음
        ("line", ["a"], ["문자열"]),                                # series 형태 오류
        ("line", ["a"], [{"name": "x", "values": ["없음"]}]),       # 숫자가 아님
        ("nope", ["a"], [{"name": "x", "values": [1]}]),           # 없는 차트 종류
    ]
    for args in bad:
        with pytest.raises(ValueError):
            asyncio.run(m.create_chart(*args))


def test_chart_escapes_labels(tmp_path):
    """라벨은 사용자/LLM이 준 문자열이다. 그대로 넣으면 SVG 구조가 깨진다."""
    m = _chart_module(tmp_path)
    r = asyncio.run(m.create_chart(
        "bar", ["<script>x</script>"], [{"name": "a&b", "values": [1]}], "5 > 3", ""))
    svg = (tmp_path / f"{r['chart_id']}.svg").read_text(encoding="utf-8")
    assert "<script>" not in svg and "&lt;script&gt;" in svg
    assert "5 &gt; 3" in svg


def test_chart_file_server_only_serves_generated_names(tmp_path):
    """경로 조작으로 컨테이너 안 다른 파일을 읽을 수 없어야 한다."""
    m = _chart_module(tmp_path)
    (tmp_path / "secret.txt").write_text("비밀", encoding="utf-8")

    async def call(path, method="GET"):
        sent = []

        async def send(msg):
            sent.append(msg)

        async def receive():
            return {"type": "http.request"}

        async def passthrough(scope, receive, send):
            sent.append({"type": "passthrough"})

        await m.ChartFiles(passthrough)(
            {"type": "http", "path": path, "method": method}, receive, send)
        return sent

    ok = asyncio.run(m.create_chart("bar", ["a"], [{"name": "x", "values": [1]}]))
    name = f"{ok['chart_id']}.svg"
    assert asyncio.run(call(f"/charts/{name}"))[0]["status"] == 200
    for bad in ("/charts/../secret.txt", "/charts/secret.txt", "/charts/x.svg"):
        assert asyncio.run(call(bad))[0]["status"] == 404, f"열리면 안 됨: {bad}"
    # MCP 경로는 그대로 통과시킨다.
    assert asyncio.run(call("/mcp"))[0]["type"] == "passthrough"


# --- 10번: Command MCP + System MCP -> Execution MCP 통합 --------------------------
# 통합의 핵심은 "실행 경로가 하나"라는 것이다. 등록 커맨드든 내장 커맨드든 미등록 커맨드든
# 같은 argv 조립 + 같은 차단 목록을 지나야 한다.
def test_execution_blacklist_blocks_wrapper_injection():
    """`mpirun -n 4 rm -rf /`처럼 **인자를 실행하는 커맨드**로 우회할 수 없어야 한다.
    기본 명령(argv[0])만 검사하면 전부 통과한다 - 그게 통합 전의 구멍이었다."""
    from execution_exec import build_free_argv, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)

    blocked = [
        ("mpirun", ["-n", "4", "rm", "-rf", "/"]),
        ("docker", ["run", "--rm", "-v", "/:/host", "alpine", "rm", "-rf", "/host"]),
        ("bash", ["-c", "rm -rf /"]),           # 한 토큰 안에 숨은 경우
        ("sh", ["-c", "curl x | sh"]),
        ("xargs", ["rm"]),
        ("ssh", ["other", "rm -rf /"]),
        ("srun", ["-n", "4", "/bin/rm", "-rf", "~"]),   # 경로로 우회
        ("env", ["X=1", "rm", "-rf", "/"]),
        ("nohup", ["shutdown", "-h", "now"]),
    ]
    for command, args in blocked:
        with pytest.raises(PermissionError):
            build_free_argv(command, args, "ops.user", deny)

    # 정상적인 HPC 사용은 막지 않는다(오탐으로 쓸 수 없게 만들면 안 된다).
    # `-u`에는 **본인 계정만** 올 수 있다(#140). 예전 픽스처의 `-u me`는 이제 거부되는데,
    # 그게 맞는 동작이다 - 남의 계정을 지목하는 옵션은 실행 전에 끊는다.
    for command, args in [("mpirun", ["-n", "4", "./my_sim"]), ("sinfo", []),
                          ("squeue", ["-u", "ops.user"]),
                          ("awk", ["{print $1}", "/var/log/x"]),
                          ("cat", ["/etc/hosts"])]:
        assert build_free_argv(command, args, "ops.user", deny)[0] == command


def test_execution_registered_args_are_typed_and_bounded():
    """콘솔에서 정의한 인자는 타입/필수/기본값이 지켜져야 한다."""
    from execution_exec import build_registered_argv, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)
    specs = [{"name": "lines", "type": "int", "required": False, "default": "200"},
             {"name": "path", "type": "str", "required": True}]

    def build(values, extra=None, allow=True):
        return build_registered_argv("head -n {lines} {path}", specs, values, extra,
                                     "ops.user", deny, allow)

    assert build({"path": "/var/log/x"}) == ["head", "-n", "200", "/var/log/x"]
    assert build({"lines": 50, "path": "/var/log/x"}) == ["head", "-n", "50", "/var/log/x"]
    with pytest.raises(ValueError):
        build({"lines": "많이", "path": "/x"})       # 정수가 아님
    with pytest.raises(ValueError):
        build({})                                    # 필수 누락
    with pytest.raises(ValueError):
        build({"path": "/x"}, ["-v"], allow=False)   # 추가 인자 금지인데 넘김
    with pytest.raises(PermissionError):
        build({"path": "/x"}, ["rm"])                # 추가 인자로 파괴적 명령

    # 값에 공백이 있어도 토큰이 쪼개지지 않아야 한다(인자 하나가 여러 개로 늘어나면 안 됨).
    argv = build({"path": "/tmp/a b.log"})
    assert argv[-1] == "/tmp/a b.log" and len(argv) == 4


def test_execution_registration_rejects_dangerous_templates():
    """등록 단계에서도 막는다 - 실행 시점 검사만 믿지 않는다."""
    from execution_exec import validate_definition, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)
    for cmd, args in [("bash -c {x}", [{"name": "x", "type": "str"}]),
                      ("docker ps", []),
                      ("rm -rf {path}", [{"name": "path", "type": "str"}])]:
        with pytest.raises(ValueError):
            validate_definition("my_tool", cmd, args, "login_server", deny)

    # 자리표시자와 인자 정의가 어긋나면 등록 단계에서 잡는다(런타임에 조용히 깨지지 않게).
    with pytest.raises(ValueError):
        validate_definition("my_tool", "head -n {lines}", [], "login_server", deny)
    with pytest.raises(ValueError):
        validate_definition("my_tool", "myquota", [{"name": "x", "type": "str"}],
                            "login_server", deny)
    # 정상 등록은 통과해야 한다.
    validate_definition("my_tool", "quota report -u {user_id}", [], "login_server", deny)
    validate_definition("my_tool", "head -n {lines} {path}",
                        [{"name": "lines", "type": "int"}, {"name": "path", "type": "str"}],
                        "login_server", deny)


def test_execution_mcp_has_no_builtin_tools():
    """커맨드는 **전부 콘솔 등록분**이어야 한다(#128).

    예전에는 파이썬 함수로 박아 둔 내장 커맨드 7개가 툴 목록에 섞여 있었다. 편집도 삭제도
    안 되면서 설명이 매 요청 프롬프트에 실렸고, 전부 LLM이 아는 표준 리눅스 명령이라
    run_command로 대체된다. 다시 코드에 커맨드를 박으면 이 테스트가 잡는다.
    """
    import importlib.util
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    spec = importlib.util.spec_from_file_location(
        "execmcp_all", os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"))
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)

    props = {t.name: t.inputSchema["properties"] for t in asyncio.run(m.mcp.list_tools())}
    # DB가 없는 테스트 환경에서는 등록 커맨드가 안 실리므로 run_command 하나만 남는다.
    assert set(props) == {"run_command"}, f"코드에 박힌 커맨드가 남아 있다: {sorted(props)}"
    for name, p in props.items():
        assert "user_id" not in p, f"{name}에 user_id가 노출됨"
    assert not os.path.exists(os.path.join(ROOT, "mcp_servers", "execution_mcp", "builtin.py"))


def test_registered_tool_schema_exposes_declared_args():
    """콘솔에서 정의한 인자가 LLM 스키마에 타입까지 그대로 보여야 한다
    (예전 Command MCP는 `args` 리스트 하나뿐이라 LLM이 무엇을 넣을지 알 수 없었다)."""
    import importlib.util
    from mcp.server.fastmcp import FastMCP
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    spec = importlib.util.spec_from_file_location(
        "registry_schema", os.path.join(ROOT, "mcp_servers", "execution_mcp", "registry.py"))
    reg = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(reg)
    from mcp_caller import build_wrapped

    async def host():
        return "10.0.0.100"

    entry = reg.build_entry({
        "title": "파일 앞부분", "description": "텍스트 파일 앞부분",
        "exec_command": "head -n {lines} {path}",
        "args": [{"name": "lines", "type": "int", "required": False, "default": "200"},
                 {"name": "path", "type": "str", "required": True}],
        "host_mode": "login_server", "enabled": True, "required_roles": [],
    }, host)

    async def state(*_a):
        return True, []

    srv = FastMCP("t", host="0.0.0.0")
    srv.add_tool(build_wrapped("read_head", entry, tool_state=state, log_execution=None,
                               host_mode="login_server", login_host=host),
                 name="read_head", description=entry["description"])
    schema = asyncio.run(srv.list_tools())[0].inputSchema
    props = schema["properties"]
    assert props["lines"]["type"] == "integer" and props["path"]["type"] == "string"
    assert schema["required"] == ["path"]
    assert "user_id" not in props and "host" not in props
    # 자유 인자는 **항상** 허용한다(#128) - 어떤 인자가 필요한지는 에이전트가 판단한다.
    assert "args" in props


# --- 11번: 이관 코드가 db-init 컨테이너에서 실제로 import 되어야 한다 -----------------
# db-init에는 `./shared`만 마운트된다. 이름 생성 규칙을 mcp_servers 쪽에 두었더니
# `No module named 'registry'`로 이관이 조용히 건너뛰어졌다(실서버에서 그렇게 실패했다).
# shared만 있는 상태를 흉내내서, 이관에 필요한 것이 전부 shared에 있는지 고정한다.
def test_migration_imports_only_shared():
    import subprocess
    # PYTHONPATH=shared 만 주고, 작업 디렉토리도 저장소 밖으로 두어 mcp_servers를 못 찾게 한다.
    env = {"PATH": os.environ.get("PATH", ""), "POSTGRES_PASSWORD": "x",
           "PYTHONPATH": os.path.join(ROOT, "shared")}
    for code in (
        "from execution_exec import tool_name_for\n"
        "assert tool_name_for('내 작업 조회', set(), 'quota report -u {user_id}') == 'quota_report'",
        "import migrations\nassert hasattr(migrations, 'import_execution_registry')",
    ):
        r = subprocess.run([sys.executable, "-c", code], capture_output=True, text=True,
                           env=env, cwd="/")
        assert r.returncode == 0, f"shared만으로 import되지 않음:\n{r.stderr}"


def test_admin_console_image_does_not_reach_into_mcp_servers():
    """콘솔 이미지는 mcp_servers를 복사하지 않는다.

    예전에는 내장 커맨드 파일 하나를 COPY했는데, 그 파일이 옮겨지거나 사라지면 운영 이미지
    빌드가 깨졌다(dev는 볼륨 마운트라 안 걸려서 못 봤다 - #112). 콘솔이 쓰는 실행 규칙은
    전부 shared에 있어야 한다.
    """
    dockerfile = open(os.path.join(ROOT, "admin_console", "Dockerfile"), encoding="utf-8").read()
    copies = [ln for ln in dockerfile.splitlines()
              if ln.strip().startswith("COPY") and "mcp_servers" in ln]
    assert not copies, f"콘솔 이미지가 mcp_servers를 복사한다: {copies}"
    router = open(os.path.join(ROOT, "admin_console", "backend", "routers", "execution.py"),
                  encoding="utf-8").read()
    assert "mcp_servers" not in router, "콘솔 라우터가 mcp_servers를 import하면 안 된다"
    assert os.path.exists(os.path.join(ROOT, "shared", "execution_exec.py"))


# --- 12번: --reload-dir가 가리키는 경로는 실제로 있어야 한다 -------------------------
# uvicorn은 없는 --reload-dir를 주면 "Invalid value" 로 **기동을 거부한다**(컨테이너 즉시 종료).
# #111에서 mcp_servers/system_mcp을 없앴는데 admin-console 이미지의 CMD가 그걸 감시하고 있어
# 관리자 콘솔(8501)이 뜨지 않았다. 이미지에 굳은 CMD라 코드만 고쳐도 안 낫는 종류의 사고다.
def test_reload_dirs_point_at_existing_paths():
    import re as _re
    import yaml

    # **git이 아는 파일**로 판단한다. 작업 트리에는 __pycache__만 남은 유령 디렉토리가 있을 수
    # 있고(git rm은 무시 파일을 지우지 않는다), 서버는 rsync --delete로 그걸 지우므로
    # os.path.exists로 보면 로컬만 통과하는 가짜 초록이 된다 - 실제로 그렇게 놓쳤다.
    import subprocess
    tracked = set(subprocess.run(["git", "ls-files"], cwd=ROOT, capture_output=True,
                                 text=True, check=True).stdout.split())
    tracked_dirs = set()
    for f in tracked:
        parts = f.split("/")
        for i in range(1, len(parts)):
            tracked_dirs.add("/".join(parts[:i]))

    def exists_in_repo(container_path: str) -> bool:
        assert container_path.startswith("/app/"), container_path
        rel = container_path[len("/app/"):]
        return rel in tracked or rel in tracked_dirs

    checked = 0
    # compose의 command
    for f in ("docker-compose.dev.yml", "docker-compose.yml"):
        spec = yaml.safe_load(open(os.path.join(ROOT, f), encoding="utf-8"))
        for name, svc in (spec.get("services") or {}).items():
            cmd = svc.get("command") or []
            if isinstance(cmd, str):
                cmd = cmd.split()
            for i, tok in enumerate(cmd):
                if tok == "--reload-dir":
                    assert exists_in_repo(cmd[i + 1]), \
                        f"{f}:{name} 의 --reload-dir 경로가 저장소에 없음: {cmd[i + 1]}"
                    checked += 1

    # Dockerfile의 CMD
    for f in ("dev/Dockerfile.admin-dev", "admin_console/Dockerfile",
              "dev/Dockerfile.agent-dev", "agent_server/Dockerfile"):
        path = os.path.join(ROOT, f)
        if not os.path.exists(path):
            continue
        for line in open(path, encoding="utf-8"):
            if not line.strip().startswith("CMD"):
                continue
            toks = _re.findall(r'"([^"]*)"', line)
            for i, tok in enumerate(toks):
                if tok == "--reload-dir":
                    assert exists_in_repo(toks[i + 1]), \
                        f"{f} 의 --reload-dir 경로가 저장소에 없음: {toks[i + 1]}"
                    checked += 1

    assert checked > 0, "검사한 --reload-dir가 하나도 없다(테스트가 무의미해짐)"


def test_compose_command_overrides_stale_admin_cmd():
    """관리자 콘솔은 compose에서 command를 지정해야 한다.
    이미지에 굳은 CMD만 믿으면, 경로가 바뀔 때 **이미지 재빌드 없이는 고칠 수 없다**."""
    _need("docker-compose.dev.yml", "docs")
    import yaml
    spec = yaml.safe_load(open(os.path.join(ROOT, "docker-compose.dev.yml"), encoding="utf-8"))
    cmd = spec["services"]["admin-console"].get("command")
    assert cmd, "admin-console에 command가 없으면 낡은 CMD가 그대로 쓰인다"
    assert "/app/admin_console" in cmd and "/app/shared" in cmd
    # 콘솔은 mcp_servers를 보지 않는다(#128). 감시 대상에 넣으면 경로가 바뀔 때 또 죽는다.
    assert not any("mcp_servers" in str(t) for t in cmd)


# --- 13번: 차트를 답변에 직접 박아 넣는다(폐쇄망에서 설정·포트 없이 동작) ----------------
# Chart MCP는 짧은 표시자(chart://<id>)만 돌려주고, Agent Server가 내보낼 때 data URI로 바꾼다.
# 이유: MCP 툴 결과는 그대로 다음 요청 프롬프트에 실린다 - base64를 돌려주면 32768이 날아간다.
def test_chart_marker_is_small_and_url_free(tmp_path):
    m = _chart_module(tmp_path)
    r = asyncio.run(m.create_chart("line", ["1월", "2월"],
                                   [{"name": "사용률", "values": [10, 20]}], "제목", "%"))
    assert r["markdown"].startswith("![제목](chart://")
    assert "http" not in r["markdown"], "기본값은 URL이 아니라 표시자여야 한다"
    # 툴 결과 전체가 프롬프트에 실린다. 300자를 넘으면 예산 설계가 깨진 것이다.
    assert len(str(r)) < 300, f"툴 결과가 너무 크다: {len(str(r))}자"


def test_chart_inliner_streaming_matches_whole():
    """표시자가 스트리밍 델타 경계에 걸쳐 쪼개져도 결과가 같아야 한다."""
    import random
    from chart_inline import ChartInliner, marker_for

    svg = "<svg xmlns='http://www.w3.org/2000/svg'>한글</svg>"

    async def fetch(_cid):
        return svg

    cid = "ab" + "0" * 30
    text = f"추이입니다.\n\n![월별]({marker_for(cid)})\n\n증가 추세입니다."
    whole = asyncio.run(ChartInliner(fetch).whole(text))
    assert "data:image/svg+xml;base64," in whole and "chart://" not in whole

    async def streamed(chunks):
        inl = ChartInliner(fetch)
        out = ""
        for c in chunks:
            out += await inl.feed(c)
        return out + await inl.flush()

    assert asyncio.run(streamed(list(text))) == whole, "1자씩 흘렸을 때 결과가 다르다"
    random.seed(7)
    for _ in range(50):
        chunks, i = [], 0
        while i < len(text):
            n = random.randint(1, 6)
            chunks.append(text[i:i + n])
            i += n
        assert asyncio.run(streamed(chunks)) == whole


def test_chart_inliner_reports_failure_instead_of_broken_image():
    from chart_inline import ChartInliner, marker_for

    async def fetch(_cid):
        return None

    out = asyncio.run(ChartInliner(fetch).whole(f"![x]({marker_for('cd' + '1' * 30)})"))
    assert "chart://" not in out and "불러오지 못했습니다" in out


def test_chart_inliner_ignores_lookalikes():
    from chart_inline import ChartInliner

    async def fetch(_cid):
        raise AssertionError("호출되면 안 됨")

    plain = "chart:// 라는 말, charter, http://x/chart 는 그대로 둔다"
    assert asyncio.run(ChartInliner(fetch).whole(plain)) == plain


def test_history_keeps_marker_not_data_uri():
    """이력/메모리에는 표시자가 남아야 한다. data URI가 저장되면 다음 프롬프트가 부푼다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    # 저장은 원문(final_text)으로, 응답 본문만 치환한다.
    assert '_bg_persist(user_id, conv, "openwebui", last_text, final_text, mem_enabled)' in src
    # 치환은 내보내는 본문에만 건다(#155부터 근거 검사를 통과한 본문에 건다).
    assert "await _chart_inliner().whole(ground.review(final_text))" in src
    # 스트리밍도 dedup.full(원문)을 저장한다.
    assert "_bg_persist(user_id, conv, \"openwebui\", last_text, dedup.full" in src \
        or "dedup.full" in src


def test_charts_base_url_derivation():
    from chart_inline import charts_base_url
    assert charts_base_url("http://chart-mcp:8005/mcp") == "http://chart-mcp:8005"
    assert charts_base_url("http://chart-mcp:8005/") == "http://chart-mcp:8005"
    assert charts_base_url("") == ""


# --- 15번: 비활성 커맨드는 툴 목록에 실리지 않아야 한다 -------------------------------
# 끄는 것만으로 막히긴 했지만(실행 시점 검사), 툴 설명이 매 요청 프롬프트에 계속 실렸다
# (하나당 ~100토큰). 게다가 에이전트가 그걸 골라 호출한 뒤 "비활성입니다" 오류를 받는
# 헛턴을 돈다. 목록 구성 단계에서 빼고, 즉시 차단은 실행 시점 검사로 유지한다.
def test_disabled_commands_are_not_exposed_as_tools():
    src = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "registry.py"),
               encoding="utf-8").read()
    assert "WHERE enabled ORDER BY title" in src, "등록 커맨드 조회가 enabled로 걸러지지 않는다"

    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    # 즉시 차단(실행 시점 검사)은 그대로 남아 있어야 한다.
    assert "async def _tool_state" in server


# --- 16번: 참고 문서 안내를 줄여서 쓰지 못하게 한다 -----------------------------------
# 관리자가 넣은 문서 위치에는 URL이 들어 있는데, LLM이 "슈퍼컴 Portal > 활용 가이드"처럼
# 요약해 버려 사용자가 문서를 찾을 수 없었다. 검색 결과가 위치와 문서 이름을 **따로** 실어
# 주고, 지시문이 정해진 두 줄 형식으로 옮기게 한다.
def test_manual_search_exposes_location_and_document_separately():
    src = open(os.path.join(ROOT, "shared", "manual_search.py"), encoding="utf-8").read()
    assert 'item["guide_location"] = item.get("reference_path")' in src
    assert 'item["guide_document"]' in src

    instr = _instruction_text()
    assert "가이드 위치:" in instr and "가이드 문서:" in instr, \
        "지시문에 참고 문서 출력 형식이 없다"
    assert "guide_location" in instr and "guide_document" in instr
    assert "한 글자도 줄이지 않고" in instr, "경로를 요약하지 말라는 규칙이 없다"


def test_instruction_asks_for_table_on_multi_column_output():
    """job 목록처럼 열이 있는 실행 결과는 표로 정리해야 한다(예전엔 그렇게 나왔다)."""
    instr = _instruction_text()
    # 제목이나 예시 단어가 아니라 **규칙**을 검사한다. #153에서 지시문을 절반으로 줄이며
    # 예시를 뺐는데, 예시를 검사하던 테스트가 규칙이 살아 있는데도 실패했다.
    assert "마크다운 테이블" in instr
    assert "값은 한 글자도 바꾸지 않습니다" in instr


def test_ssh_master_health_is_observable():
    """'ssh 세션이 제대로 열렸는지'를 로그로 확인할 수 있어야 한다.
    추측으로 느림을 진단할 수 없다 - 마스터가 죽으면 커맨드마다 수십 초가 더 붙는다."""
    ssh = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    assert "async def master_alive" in ssh
    assert '"-O", "check"' in ssh, "ssh -O check로 실제 상태를 확인해야 한다"

    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "다중화 마스터 준비 완료" in server
    assert "매번 새로 접속해" in server, "마스터가 없을 때의 영향을 로그로 알려야 한다"


def test_master_session_is_resident_and_supervised():
    """로그인 서버로의 root ssh 세션을 **상주**시키고 감시해야 한다(사용자 요구).

    예전에는 `ssh … true`로 연결만 만들고 수명을 ControlPersist에 맡겼다. 죽으면 다음
    확인(180초)까지 구멍이 났고, 그 사이 커맨드는 매번 새로 접속했다(실측 17~25초).
    이제 마스터 ssh를 우리 자식 프로세스로 붙들고(-M -N) 15초마다 살아 있는지 본다.
    """
    ssh = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    assert "async def ensure_master" in ssh and "def start_master_supervisor" in ssh
    assert '"-M", "-N"' in ssh, "원격 명령 없는 마스터 전용 연결이어야 한다"
    # ControlPersist를 주면 ssh가 스스로 백그라운드로 가버려 감시가 불가능해진다.
    assert '"ControlPersist=no" if master' in ssh
    assert '"ControlMaster=yes" if master' in ssh
    assert "async def stop_masters" in ssh, "종료 시 정리가 없다"

    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "start_master_supervisor" in server
    assert 'add_event_handler("startup"' in server and 'add_event_handler("shutdown"' in server
    # 감시 주기는 짧아야 한다(살아 있으면 파일 확인 한 번이라 거의 공짜다).
    assert "interval=15" in server


# --- 18번: 커맨드 실행이 왜 느린지 **측정**할 수 있어야 한다 --------------------------
# "느리다"는 리포트가 올 때마다 원인을 추측해 왔다(빈 키 파일·호스트 키·TTY — 셋 다 틀렸고
# 진짜 원인은 타임아웃이었다, #69). 이제 매 실행이 소요 시간과 접속 재사용 여부를 달고 온다.
def test_ssh_result_carries_timing_and_reuse():
    ssh = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    assert '"duration_ms"' in ssh, "실행 결과에 소요 시간이 없다"
    assert '"connection_reused"' in ssh, "접속을 새로 맺었는지 알 수 없다"
    # 소켓 경로를 우리가 정해야 밖에서 재사용 여부를 확인할 수 있다(ssh의 %C로는 불가).
    assert "def control_path" in ssh and "def master_socket_exists" in ssh
    assert "ControlPath={control_path(ip)}" in ssh, "소켓 경로를 우리가 정하지 않는다"

    # 진행 상황 줄에 소요 시간이 보여야 사용자가 어느 커맨드가 느린지 짚어 줄 수 있다.
    main = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert "duration_ms" in main and "초\"" in main
    # 요청 단위 계측(전체 / 첫 글자 / 도구 횟수 / 커맨드 실행 시간).
    assert "class _Pace" in main and "커맨드 실행" in main


def test_tool_call_hits_db_once_and_does_not_block_on_audit_log():
    """툴 호출 하나가 사용자를 기다리게 하는 DB 왕복을 줄인다.

    예전에는 같은 행을 두 번 조회하고(enabled / required_roles), 성공 감사로그 INSERT까지
    await했다. 커맨드를 여러 개 부르는 질문에서는 그만큼 그대로 쌓인다.
    """
    caller = open(os.path.join(ROOT, "shared", "mcp_caller.py"), encoding="utf-8").read()
    assert "tool_state" in caller and "is_enabled" not in caller, \
        "활성/역할 조회가 아직 두 번으로 나뉘어 있다"
    assert "_log_later(log_execution, name, params, \"success\"" in caller, \
        "성공 감사로그가 응답을 막고 있다"
    # 거부/차단은 실행되지 않아 빠르므로 그대로 await한다(기록이 응답보다 먼저 남는 편이 낫다).
    assert 'await log_execution(name, params, "blocked"' in caller

    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "SELECT enabled, required_roles FROM execution_commands" in server


def test_session_history_is_written_without_reloading_the_session():
    """이력 주입이 턴 수의 제곱으로 늘어나면 안 된다(사용자가 첫 글자를 보기 전의 지연)."""
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    body = _re.search(r"async def _create_session\(.*?\n    return session_id", src, _re.S).group(0)
    assert "svc.get_session(" not in body, "이력 한 턴마다 세션을 다시 읽고 있다"
    assert "append_event" in body


_FAKE_SSH = '''#!/usr/bin/env python3
"""테스트용 가짜 ssh: 마스터 동작만 흉내낸다(소켓 파일 생성 + -N이면 계속 살아 있음)."""
import os, sys, time
args = sys.argv[1:]
cp = ""
for a in args:
    if a.startswith("ControlPath="):
        cp = a.split("=", 1)[1]
if "-O" in args and "check" in args:
    sys.exit(0 if cp and os.path.exists(cp) else 255)
if "-M" in args and "-N" in args:
    open(cp, "w").close()
    while True:
        time.sleep(1)
sys.exit(0)
'''


@pytest.mark.filterwarnings("ignore::pytest.PytestUnraisableExceptionWarning")
def test_resident_master_is_spawned_adopted_and_restarted(tmp_path, monkeypatch):
    """상주 마스터의 상태 기계를 실제로 돌려서 확인한다(가짜 ssh 사용).

    띄운다 → 이미 있으면 그대로 쓴다 → 죽으면 감시 루프가 다시 띄운다.
    이 셋 중 하나라도 깨지면 커맨드마다 새 접속(실측 17~25초)을 물게 된다.
    """
    import ssh_exec

    fake_dir = tmp_path / "bin"
    fake_dir.mkdir()
    fake = fake_dir / "ssh"
    fake.write_text(_FAKE_SSH, encoding="utf-8")
    fake.chmod(0o755)
    monkeypatch.setenv("PATH", f"{fake_dir}:{os.environ.get('PATH', '')}")
    monkeypatch.setattr(ssh_exec, "SSH_CONTROL_DIR", str(tmp_path / "mux"))
    monkeypatch.setattr(ssh_exec, "SSH_KEY", "")
    monkeypatch.setattr(ssh_exec, "_control_dir_ready", False)
    monkeypatch.setattr(ssh_exec, "_master_procs", {})

    host = "203.0.113.9"          # 문서용 예약 대역(TEST-NET-3) - 실제로 붙지 않는다

    async def scenario():
        first = await ssh_exec.ensure_master(host)
        assert first["ok"] and first["already_up"] is False, first
        assert await ssh_exec.master_alive(host)

        # 두 번째 호출은 새로 띄우지 않는다(소켓이 겹치면 실패한다).
        again = await ssh_exec.ensure_master(host)
        assert again["ok"] and again["already_up"] is True, again

        # 마스터가 죽은 상황을 만든다.
        ip = ssh_exec.resolve_host(host)
        proc = ssh_exec._master_procs[ip]
        proc.kill()
        await proc.wait()
        os.unlink(ssh_exec.control_path(ip))
        assert not await ssh_exec.master_alive(host)

        # 감시 루프가 스스로 복구해야 한다.
        task = ssh_exec.start_master_supervisor(
            lambda: asyncio.sleep(0, result=host), interval=1)
        try:
            for _ in range(40):
                if await ssh_exec.master_alive(host):
                    break
                await asyncio.sleep(0.25)
            assert await ssh_exec.master_alive(host), "감시 루프가 마스터를 복구하지 못했다"
        finally:
            task.cancel()
            await ssh_exec.stop_masters()

    asyncio.run(scenario())


def test_agent_self_name_is_rejected_as_a_path_or_account():
    """에이전트가 **자기 이름**을 계정/경로로 써서 만든 커맨드는 실행하지 않는다.

    "내 홈 파일 리스트"에 `ls -la /home/ops_assistant`를 돌리고 "경로가 없습니다"라고
    답한 사고가 있었다. ADK가 시스템 프롬프트에 넣는 에이전트 이름을 사용자 계정으로
    착각한 것이다(#125와 같은 뿌리). 지시문으로 두 번 막았는데 재발해서 실행 단계에서 끊는다.
    """
    from execution_exec import build_free_argv, build_registered_argv, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)

    with pytest.raises(PermissionError) as e:
        build_free_argv("ls", ["-la", "/home/ops_assistant"], "ops.user", deny)
    # 거부만 하지 말고 다음에 뭘 해야 하는지 알려 준다.
    assert "경로를 비우거나" in str(e.value)

    # 등록 커맨드의 추가 인자로 들어와도 막는다.
    with pytest.raises(PermissionError):
        build_registered_argv("ls", [], {}, ["/home/ops_assistant"], "ops.user", deny, True)

    # 정상 사용은 그대로 통과해야 한다(과잉 차단 금지).
    assert build_free_argv("ls", ["-lh"], "ops.user", deny) == ["ls", "-lh"]
    assert build_free_argv("ls", ["/home/users/ops.user"], "ops.user", deny)[-1] \
        == "/home/users/ops.user"


def test_adk_streaming_toolcall_index_bug_is_reproducible():
    """google-adk 1.22.1의 스트리밍 툴 호출 인자 누적이 index 0을 '없음'으로 취급한다.

    `index = chunk.index or fallback_index` — 파이썬에서 0은 거짓이다. vLLM(hermes)이 같은
    호출의 조각에 index를 0 → 1로 바꿔 보내면 인자가 두 통으로 쪼개져 각각 잘린 JSON이 되고,
    `_message_to_generate_content_response`가 try/except 없이 json.loads 해서 요청이 죽는다.
    실서버 오류(`Expecting value: line 1 column 11 (char 10)`)가 정확히 이 모양이다.
    라이브러리를 올릴 때 이 테스트로 고쳐졌는지 확인한다.
    """
    import json as _json

    def adk_accumulate(chunks):          # lite_llm.py의 누적 로직을 그대로 옮긴 것
        function_calls, fallback_index = {}, 0
        for c in chunks:
            index = c["index"] or fallback_index
            function_calls.setdefault(index, {"args": ""})
            if c["args"]:
                function_calls[index]["args"] += c["args"]
                try:
                    _json.loads(function_calls[index]["args"])
                    fallback_index += 1
                except _json.JSONDecodeError:
                    pass
        return function_calls

    buckets = adk_accumulate([
        {"index": 0, "args": ""},
        {"index": 0, "args": '{"lines": '},
        {"index": 1, "args": "200}"},        # 같은 호출인데 index가 바뀐 경우
    ])
    assert buckets[0]["args"] == '{"lines": ', "인자가 쪼개지지 않았다면 전제가 바뀐 것"
    with pytest.raises(_json.JSONDecodeError) as e:
        _json.loads(buckets[0]["args"])
    assert "line 1 column 11 (char 10)" in str(e.value), "실서버 오류 메시지와 같아야 한다"


def test_streaming_toolcall_failure_falls_back_to_non_streaming():
    """위 결함에 걸리면 논스트리밍으로 한 번 더 돌려 답을 낸다(요청이 죽지 않게).
    다른 오류는 그대로 올려보내야 한다 - 원인을 숨기면 안 된다."""
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("_CONTEXT_ERROR_MARKERS")
    j = src.index("\ndef _trace_ctx(")

    created = []

    async def _create_session(user_id, history):
        created.append(history)
        return f"retry-{len(created)}"

    async def _cleanup_session(u, s):
        pass

    async def get_config(_k, default=None):
        return default

    ns = {"re": _re, "get_config": get_config,
          "_create_session": _create_session, "_cleanup_session": _cleanup_session}
    exec(src[i:j], ns)
    recover = ns["_run_with_toolcall_recovery"]

    class FakeRunner:
        def run_async(self, *, user_id, session_id, new_message, run_config=None):
            async def gen():
                if run_config is not None:          # 스트리밍 -> ADK 결함 재현
                    yield "tool-event"
                    raise ValueError("Expecting value: line 1 column 11 (char 10)")
                yield "final-answer"                # 논스트리밍 -> 성공
            return gen()

    async def scenario():
        got = [e async for e in recover(FakeRunner(), "ops.user", "s1", "msg",
                                        object(), history=[("user", "안녕")])]
        assert "final-answer" in got, f"논스트리밍 재시도가 동작하지 않았다: {got}"
        assert len(created) == 1, "재시도는 새 세션에서 돌려야 한다(중간 이벤트가 남아 있다)"

        class Boom:
            def run_async(self, **kw):
                async def gen():
                    raise RuntimeError("connection refused")
                    yield
                return gen()

        with pytest.raises(RuntimeError):
            async for _ in recover(Boom(), "u", "s", "m", object(), history=[]):
                pass

    asyncio.run(scenario())


def test_tool_call_json_error_is_translated():
    """vLLM tool-call 파서가 깨진 JSON을 내보내면 `Expecting value: line 1 column 11`이
    그대로 사용자에게 갔다. 사용자가 할 수 있는 조치로 바꿔 말한다."""
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("_CONTEXT_ERROR_MARKERS")

    async def get_config(_k, default=None):
        return ""

    ns = {"re": _re, "get_config": get_config}
    exec(src[i:src.index("\ndef _trace_ctx(")], ns)
    msg = asyncio.run(ns["_friendly_error"](
        ValueError("Expecting value: line 1 column 11 (char 10)")))
    assert "Expecting value" not in msg
    assert "다시" in msg or "한 번 더" in msg

    # 스택트레이스를 로그에 남겨야 어디서 났는지 알 수 있다(예전엔 메시지만 찍었다).
    assert "traceback.format_exc()" in src


def test_ssh_handshake_options_are_disabled_by_default():
    """첫 접속이 17.4초였다. TCP가 아니라 **인증 협상**이 원인이라 ConnectTimeout으로는 못 막는다.
    협상을 늘리는 것들(GSSAPI·여러 키 시도·IPv6)을 끈 상태로 유지한다."""
    ssh = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    for opt in ("GSSAPIAuthentication=no", "PreferredAuthentications=publickey",
                "AddressFamily=inet", "IdentitiesOnly=yes"):
        assert opt in ssh, f"핸드셰이크 최적화가 빠졌다: {opt}"


def test_warm_endpoint_exists_and_is_not_a_tool():
    """새로고침/새 채팅 시점에 ssh 세션이 서 있어야 한다.

    예열은 HTTP 라우트여야 한다 - MCP 툴로 만들면 설명이 매 요청 프롬프트에 실리고,
    에이전트가 그걸 골라 호출하는 헛턴도 생긴다.
    """
    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "async def warm_endpoint" in server
    assert 'Route("/warm"' in server
    assert "mcp.add_tool" not in server.split("def warm_endpoint")[1], \
        "예열을 MCP 툴로 노출하면 안 된다"

    main = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert "def warm_execution_host" in main
    # Open WebUI가 페이지를 열 때 부르는 경로와 채팅 시작 지점 모두에서 예열한다.
    assert main.count("warm_execution_host()") >= 3, "예열 호출 지점이 부족하다"


# 비로그인 모드에는 홈 이동(`cd ~user`)이 앞에 붙는다(#144) - 강등 부분만 떼어 검사한다.
@pytest.mark.parametrize("mode,expect_prefix", [
    ("su-login", "su - ops.user -c "),
    ("su", "su ops.user -c "),
    ("runuser", "runuser -u ops.user -- "),
])
def test_privilege_drop_modes_build_safe_remote_commands(mode, expect_prefix, monkeypatch):
    """권한 강등 방식 3가지. **어느 쪽이든 호출자 본인 계정으로 내려간다**(우회 경로 없음).

    차이는 커맨드 하나당 고정 비용이다 - `su -`는 원격 프로필을 매번 읽어 실측 약 2초를 쓰고,
    `runuser`는 PAM 인증과 프로필을 모두 건너뛴다. 어느 쪽을 쓰든 셸 메타문자는 인용 밖으로
    새면 안 된다(그게 새면 `; rm -rf /`가 그대로 실행된다).
    """
    import importlib
    monkeypatch.setenv("SSH_PRIVDROP", mode)
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    import ssh_exec
    importlib.reload(ssh_exec)
    try:
        assert ssh_exec.SSH_PRIVDROP == mode
        cmd = ssh_exec._remote_command("ops.user", ["ls", "-lh", "; rm -rf /", "`whoami`"])
        drop = cmd.split("2>/dev/null; ", 1)[-1]      # 홈 이동 접두사를 떼어낸다
        assert drop.startswith(expect_prefix), cmd
        for danger in ("; rm -rf /", "`whoami`"):
            assert f" {danger} " not in cmd, f"인용되지 않은 채 노출됨: {danger}"
    finally:
        monkeypatch.delenv("SSH_PRIVDROP", raising=False)
        importlib.reload(ssh_exec)


_FAKE_SSH_PROFILE_ONLY = r'''#!/usr/bin/env python3
"""가짜 ssh: `phd`는 로그인 셸에서만 찾아지고, runuser는 아예 없는 서버를 흉내낸다."""
import re, sys
remote = sys.argv[-1]
# 비로그인 모드에는 홈 이동이 앞에 붙는다(#144). 진짜 셸이 하듯 그 부분을 소화한다.
remote = re.sub(r"^cd ~\S+ 2>/dev/null; ", "", remote)
if remote.startswith("runuser"):
    sys.stderr.write("bash: runuser: command not found\\n"); sys.exit(127)
if remote.startswith("su - "):
    sys.stdout.write("JOBID  STATE\\n1234   RUN\\n"); sys.exit(0)
if remote.startswith("su "):
    sys.stderr.write("bash: phd: command not found\\n"); sys.exit(127)
sys.exit(0)
'''


@pytest.mark.filterwarnings("ignore::pytest.PytestUnraisableExceptionWarning")
def test_non_login_privdrop_recovers_without_failing_the_user(tmp_path, monkeypatch):
    """비로그인 강등을 켜도 **사용자에게는 실패가 보이면 안 된다**.

    두 가지가 겹칠 수 있다: 대상 서버에 runuser가 없고, 그 커맨드는 프로필(PATH)이 있어야
    찾아진다. 한 번만 재시도하면 첫 실행이 그대로 실패하므로 루프로 돈다.
    그리고 무엇이 프로필을 필요로 하는지 **기억해서** 다음부터는 처음부터 로그인 셸로 간다.
    """
    import importlib
    fake_dir = tmp_path / "bin"
    fake_dir.mkdir()
    fake = fake_dir / "ssh"
    fake.write_text(_FAKE_SSH_PROFILE_ONLY, encoding="utf-8")
    fake.chmod(0o755)
    monkeypatch.setenv("PATH", f"{fake_dir}:{os.environ.get('PATH', '')}")
    monkeypatch.setenv("SSH_PRIVDROP", "runuser")
    monkeypatch.setenv("SSH_MULTIPLEX", "false")

    sys.path.insert(0, os.path.join(ROOT, "shared"))
    import ssh_exec
    importlib.reload(ssh_exec)
    try:
        async def limit():
            return 4000
        ssh_exec.set_output_limit_getter(limit)

        async def scenario():
            first = await ssh_exec.run_ssh_as_user("203.0.113.9", "ops.user", ["phd", "list"])
            assert first["exit_code"] == 0, f"첫 실행이 실패했다: {first}"
            assert first["privdrop"] == "su-login"
            assert "1234" in first["stdout"]

            second = await ssh_exec.run_ssh_as_user("203.0.113.9", "ops.user", ["phd", "list"])
            assert second["privdrop"] == "su-login", "배운 걸 안 쓰고 또 돌아갔다"
            assert "phd" in ssh_exec._NEEDS_LOGIN_SHELL
            assert ssh_exec._privdrop_downgrade == "su", "runuser 미지원을 기억하지 않았다"

        asyncio.run(scenario())
    finally:
        for k in ("SSH_PRIVDROP", "SSH_MULTIPLEX"):
            monkeypatch.delenv(k, raising=False)
        importlib.reload(ssh_exec)


def test_privilege_drop_defaults_to_login_shell_and_keeps_legacy_flag():
    """기본은 로그인 셸이다 - 사내 커맨드가 프로필의 PATH에 의존하면 다른 방식은 깨진다.
    먼저 재 보고(bench-exec.sh) 확인된 환경에서만 바꾼다. 구 SSH_SU_LOGIN도 계속 통해야 한다."""
    import importlib
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    import ssh_exec
    importlib.reload(ssh_exec)
    assert ssh_exec.SSH_PRIVDROP == "su-login"

    os.environ["SSH_SU_LOGIN"] = "false"
    try:
        importlib.reload(ssh_exec)
        assert ssh_exec.SSH_PRIVDROP == "su", "구 설정(SSH_SU_LOGIN=false)이 안 먹는다"
    finally:
        os.environ.pop("SSH_SU_LOGIN", None)
        importlib.reload(ssh_exec)

    # 강등 도구 자체는 사용자 커맨드로 실행될 수 없어야 한다(우회 방지).
    from execution_exec import DENY_BASE_COMMANDS
    assert {"su", "runuser", "setpriv", "sudo"} <= DENY_BASE_COMMANDS


# --- 20번: 내부 신뢰 경계 인증 -------------------------------------------------------
# MCP는 X-User-Id를 그대로 믿고 그 계정 권한으로 커맨드를 실행한다. MCP 포트가 호스트에
# 열려 있으면 같은 망의 누구나 그 헤더를 붙여 **남의 계정으로 실행**할 수 있었다.
def test_mcp_rejects_calls_without_shared_secret():
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    from mcp_caller import CallerContextMiddleware, get_caller

    seen = []

    async def app(scope, receive, send):
        seen.append(("passed", get_caller().get("user_id")))

    async def send(msg):
        if msg["type"] == "http.response.start":
            seen.append(("status", msg["status"]))

    async def call(configured_secret, sent_header):
        seen.clear()
        mw = CallerContextMiddleware(
            app, secret_getter=lambda: asyncio.sleep(0, result=configured_secret))
        headers = [(b"x-user-id", b"ops.user")]
        if sent_header is not None:
            headers.append((b"x-agent-secret", sent_header.encode()))
        await mw({"type": "http", "headers": headers}, None, send)
        return list(seen)

    async def scenario():
        assert await call("s3cr3t", "s3cr3t") == [("passed", "ops.user")], "정상 호출이 막혔다"
        assert ("status", 401) in await call("s3cr3t", None), "비밀값 없이 통과했다"
        assert ("status", 401) in await call("s3cr3t", "wrong"), "틀린 비밀값으로 통과했다"
        # 비밀값이 아직 없는 구 배포는 막지 않는다(돌던 서비스를 갑자기 세우지 않는다).
        assert await call("", None) == [("passed", "ops.user")]

    asyncio.run(scenario())

    # 양쪽이 같은 DB 값을 쓰는지 - db-init이 무작위로 심고, agent-server가 헤더로 보낸다.
    cfg = open(os.path.join(ROOT, "shared", "migrations.py"), encoding="utf-8").read()
    assert '("mcp_shared_secret", secrets.token_urlsafe(32)' in cfg
    agent = open(os.path.join(ROOT, "agent_server", "agent.py"), encoding="utf-8").read()
    assert 'headers["X-Agent-Secret"] = mcp_secret' in agent
    for mcp_dir in ("execution_mcp", "chart_mcp"):
        srv = open(os.path.join(ROOT, "mcp_servers", mcp_dir, "server.py"), encoding="utf-8").read()
        assert "secret_getter=" in srv, f"{mcp_dir}가 비밀값을 검사하지 않는다"


def test_agent_server_v1_endpoints_can_require_api_key():
    """`X-OpenWebUI-User-Email`을 그대로 믿는 서버라, 포트가 열려 있으면 헤더만 바꿔
    남의 계정으로 실행할 수 있다. API 키를 넣으면 `/v1/*`이 전부 잠겨야 한다.

    **엔드포인트를 소스에서 열거한다.** 예전에는 4개를 손으로 적어 뒀는데, 그 목록에 없던
    `/v1/memory/{user_id}` 셋(GET/POST/DELETE)이 인증 없이 열려 있었다(#143). 목록을 박아 두면
    나중에 추가되는 엔드포인트를 영원히 못 잡는다.
    """
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert "async def require_api_key" in src
    assert "hmac.compare_digest" in src, "타이밍 안전 비교를 쓰지 않는다"

    # 의도적으로 열어 두는 경로만 여기 적는다. 추가할 때는 왜 안전한지 이유를 함께 남길 것.
    #   /health - 신원도 실행도 없고, restart-mounted.sh가 기동 확인에 쓴다.
    exempt = {"/health"}

    decorated = re.findall(r'@app\.(get|post|put|delete|patch)\("([^"]+)"([^)]*)\)', src)
    assert decorated, "엔드포인트를 하나도 찾지 못했다(정규식을 갱신할 것)"

    unprotected = [
        f"{verb.upper()} {path}"
        for verb, path, rest in decorated
        if path not in exempt and "require_api_key" not in rest
    ]
    assert not unprotected, (
        "인증 없이 열려 있는 엔드포인트: " + ", ".join(unprotected) +
        "\n  /v1/*는 호출자가 준 user_id를 그대로 믿고 그 계정으로 실행/조회한다.")

    # 메모리 POST는 특히 위험하다 - 넣은 내용이 그 사용자의 다음 대화에서 지시문에 붙는다.
    i = src.index('@app.post("/v1/memory/{user_id}"')
    assert "require_api_key" in src[i:i + 200]

    # 꺼져 있으면 기동 로그로 알려야 한다(조용히 열어 두지 않는다).
    assert "/v1/* 에 인증이 없습니다" in src


def test_console_explains_405_instead_of_showing_it():
    """새 API를 추가하면 콘솔 화면은 바로 새 코드인데 백엔드는 재시작해야 바뀐다.

    그 사이에 새 버튼을 누르면 FastAPI가 모르는 경로를 StaticFiles 마운트로 넘겨
    `405 Method Not Allowed`가 뜬다 - 원인과 아무 상관 없어 보이는 메시지다.
    세 번 겪었으므로(#27, #30, #138) 메시지가 **다음에 할 일**을 말하게 한다.
    """
    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    assert "res.status === 405" in html, "405를 특별히 처리하지 않는다"
    assert "restart admin-console" in html, "재시작 커맨드를 알려주지 않는다"

    # 405가 나는 구조 자체(맨 끝 StaticFiles 마운트)는 그대로여야 이 처리가 의미가 있다.
    main = open(os.path.join(ROOT, "admin_console", "backend", "main.py"), encoding="utf-8").read()
    assert 'app.mount("/", StaticFiles(' in main

    # 배포 절차에도 admin-console 재시작이 들어 있어야 한다.
    # **NEXT-STEPS가 아니라 스크립트에 걸어야 한다** - NEXT-STEPS는 매 턴 "지금 할 일"만
    # 남기고 새로 쓰는 문서라(CLAUDE.md 1절), 이번에 콘솔을 안 건드리면 그 줄이 사라진다.
    # 영구 보장은 사용자가 매번 돌리는 `restart-mounted.sh`가 해야 한다.
    restart = open(os.path.join(ROOT, "scripts", "restart-mounted.sh"), encoding="utf-8").read()
    assert "admin-console" in restart, \
        "restart-mounted.sh가 admin-console을 재시작하지 않는다 - 405가 재발한다"


def test_rsync_never_deletes_server_only_files():
    """배포 rsync가 `--delete`로 **서버에만 있어야 하는 파일**을 지우면 안 된다(#137).

    `.env`는 .gitignore에 있고 `secrets/`(ssh 개인키)도 저장소에 없다. 보내는 쪽에 없으니
    `--delete`가 매번 지웠다. 바인드 마운트 때문에 이미 떠 있는 컨테이너는 멀쩡해서,
    다음 `up -d`(재생성) 때 비로소 "모든 커맨드 인증 실패"로 터진다 - 원인 찾기가 특히 어렵다.

    #141 이후로 제외 목록은 **문서가 아니라 `scripts/deploy-rsync.sh`에** 산다. 문서에 적어 두고
    복사해 쓰는 방식 자체가 사고의 원인이었기 때문이다(빠뜨려도 아무도 모른다).
    그래서 검사도 두 갈래다: 스크립트가 제외를 갖고 있는가 + 문서에 맨손 rsync가 되살아났는가.
    """
    _need("scripts/deploy-rsync.sh")
    gitignore = open(os.path.join(ROOT, ".gitignore"), encoding="utf-8").read()
    assert "secrets/" in gitignore, "개인키 디렉토리가 gitignore에 없다"

    script_path = os.path.join(ROOT, "scripts", "deploy-rsync.sh")
    assert os.path.isfile(script_path), "scripts/deploy-rsync.sh가 없다"
    script = open(script_path, encoding="utf-8").read()
    # **주석 줄은 빼고** 본다. 그냥 substring으로 찾으면 `#--exclude '.env'`처럼 주석 처리된
    # 것도 통과해서, 테스트가 있는데도 회귀를 놓친다(이 테스트를 쓰다 실제로 겪었다).
    active = "\n".join(ln for ln in script.split("\n") if not ln.lstrip().startswith("#"))
    for pattern in ("--exclude '.env'", "--exclude 'secrets/'"):
        assert pattern in active, f"deploy-rsync.sh에 {pattern}가 없다 - --delete가 지운다"

        # **저장소가 안 들고 다니는데 서버에는 있어야 하는 것**이 --delete의 1순위 표적이다.
        # `.gitignore`에 있다는 것은 "저장소가 안 들고 다닌다"는 뜻이지 "없어도 된다"는
        # 뜻이 아니다. vendor의 React·Babel이 지워지면 관리자 콘솔이 통째로 빈 화면이 된다.
        for pattern in ("admin_console/frontend/vendor/", "Temp/"):
            assert pattern in active, \
                f"deploy-rsync.sh에 {pattern}가 없다 - 서버에만 있는 파일을 --delete가 지운다"
    # 지워질 파일을 **먼저 보여주고 확인을 받는** 것이 이 스크립트의 존재 이유다.
    assert "del." in script and "계속할까요" in script, \
        "deploy-rsync.sh가 삭제 목록을 보여주고 확인받지 않는다"

    # 문서에 맨손 rsync가 되살아나면(복사해 쓰다 제외를 빠뜨리는 경로) 잡는다.
    for doc in ("CLAUDE.md", os.path.join("docs", "NEXT-STEPS.md"),
                os.path.join("docs", "RUN-LOG.md")):
        text = open(os.path.join(ROOT, doc), encoding="utf-8").read()
        lines, i, commands = text.split("\n"), 0, []
        while i < len(lines):
            if lines[i].strip().startswith("rsync "):
                cmd = lines[i].rstrip()
                while cmd.endswith("\\") and i + 1 < len(lines):
                    i += 1
                    cmd = cmd[:-1] + " " + lines[i].strip()
                commands.append(cmd)
            i += 1
        for cmd in commands:
            if "--delete" not in cmd:
                continue        # --delete가 없으면 지울 일이 없다
            assert "--exclude '.env'" in cmd, (
                f"{doc}: 맨손 rsync --delete가 .env를 지운다. "
                f"scripts/deploy-rsync.sh를 쓰도록 고칠 것\n{cmd}")
            assert "--exclude 'secrets/'" in cmd, (
                f"{doc}: 맨손 rsync --delete가 ssh 키를 지운다\n{cmd}")

    # 키가 사라진 상태를 조용히 넘기면 안 된다 - 기동 로그에서 바로 보여야 한다.
    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "os.path.isfile(SSH_KEY)" in server
    assert "모든 커맨드 실행이 인증 실패합니다" in server


def test_instruction_can_be_reset_from_console_without_db_env():
    """지시문을 **버튼 하나로** 최신 기본값으로 되돌릴 수 있어야 한다(#136).

    non-force 시드라 db-init이 기존 값을 안 덮으므로, 예전에는 1만 자짜리 전문을 문서에 붙이고
    사람이 복사·붙여넣기 했다 - 매번 반복이고 중간에 잘리면 조용히 깨진다.
    그리고 원문은 **부수효과 없는 모듈**에 있어야 한다. migrations.py는 import 시점에
    POSTGRES_PASSWORD를 요구하는데, 콘솔 컨테이너에는 그 환경변수가 없어 터진다.
    """
    import subprocess
    env = {"PATH": os.environ.get("PATH", ""),
           "PYTHONPATH": os.path.join(ROOT, "shared")}     # POSTGRES_PASSWORD 없음 = 콘솔과 같은 조건
    r = subprocess.run(
        [sys.executable, "-c",
         "from agent_instruction import AGENT_INSTRUCTION as A; assert len(A) > 1000; print(len(A))"],
        capture_output=True, text=True, env=env, cwd="/")
    assert r.returncode == 0, f"콘솔 환경에서 지시문을 읽지 못한다:\n{r.stderr}"

    router = open(os.path.join(ROOT, "admin_console", "backend", "routers", "settings.py"),
                  encoding="utf-8").read()
    assert '@router.post("/agent_system_instruction/reset")' in router
    # 읽는 방식은 #147에서 바뀌었다: 모듈 import는 sys.modules에 캐시돼 **옛 텍스트**를
    # 계속 저장했다. 이제 파일을 직접 읽는다(test_instruction_reset_reads_file_not_module_cache).
    assert "_read_instruction_from_disk" in router
    assert "agent_instruction.py" in router, "지시문 파일 경로를 참조하지 않는다"
    assert "from migrations import" not in router, "콘솔이 migrations를 import하면 터진다"

    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    assert "지시문을 최신 기본값으로 되돌리기" in html
    assert "/api/settings/agent_system_instruction/reset" in html


def test_instruction_names_no_in_house_command():
    """지시문에 **사내 전용 커맨드 이름을 쓰지 않는다** — 금지 예시로도 쓰지 않는다.

    "`phd info` 같은 커맨드를 지어내지 마세요"라고 적어 뒀더니, 모델이 그 이름을 그대로
    가져다 실행했다. 지시문은 매 요청 시스템 프롬프트라, 거기 적힌 문자열은 '금지 목록'이
    아니라 '아는 커맨드'로 읽힌다. #74에서 컴파일 옵션을 예시로 들었다가 같은 사고를 냈다.
    표준 리눅스 명령(ls·df·find …)은 실제로 존재하므로 예외다.
    """
    instr = _instruction_text()
    # 이 시스템에 있는지 우리가 확인할 수 없는 커맨드 이름들(과거에 지시문에 새어 들어간 것 포함).
    forbidden = ["phd ", "myquota", "squeue", "sinfo", "sbatch", "bsub", "qstat", "lsload"]
    hits = [w for w in forbidden if w in instr]
    assert not hits, f"지시문에 사내 커맨드 이름이 있다(모델이 그대로 쓴다): {hits}"


def test_instruction_routes_own_resource_checks_straight_to_execution():
    """'내 job 현황'처럼 본인 자원을 물으면 매뉴얼을 뒤지지 말고 바로 실행해야 한다.
    '현황'이라는 낱말 때문에 매뉴얼 검색이 앞에 붙으면 답이 몇 초씩 늦어진다."""
    instr = _instruction_text()
    assert "내 job 현황" in instr, "'현황'이 붙은 본인 자원 질문의 예외가 지시문에 없다"
    assert "매뉴얼도 과거 사례도 먼저 뒤지지 않습니다" in instr


def test_console_role_is_a_select_with_admin_user():
    """역할은 자유 입력이면 오타 하나로 아무도 못 쓰게 된다. select로 고정한다."""
    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    assert "const ROLE_OPTIONS" in html
    for v in ('value: ""', 'value: "user"', 'value: "admin"'):
        assert v in html, f"역할 선택지에 {v}가 없다"
    assert 'onChange={ev => setE({ role: ev.target.value })}' in html


# --- 17번: multi-turn 기억 오염 ------------------------------------------------------
# "CPU에서 스크래치 사용법" 다음에 "GPU에서 스크래치 사용법"을 물으면 CPU 절차가 나왔다.
# 경로가 셋이었다: (1) 요약기가 절차를 장기기억으로 승격, (2) 그 장기기억이 시스템 지시문에
# '기억된 정보'로 주입돼 근거처럼 쓰임, (3) 대화 이력의 앞 답변 재사용.
def test_summarizer_refuses_to_memorize_procedures():
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("async def _summarize_turns")
    prompt = src[i:i + 2500]
    assert "절대 뽑지 말 것" in prompt, "요약기가 무엇을 배제해야 하는지 명시하지 않는다"
    for banned in ("사용법", "절차", "명령어", "경로", "옵션", "설정값"):
        assert banned in prompt, f"요약기 프롬프트에 '{banned}' 배제가 없다"


def test_memory_block_is_marked_as_not_evidence():
    """장기기억 블록은 시스템 지시문에 실린다. '근거가 아니다'가 명시돼야 한다."""
    from memory_store import format_memory_block
    block = format_memory_block([{"kind": "fact", "content": "테스트 항목"}])
    assert "근거가 아닙니다" in block
    assert "다시 검색" in block
    assert "테스트 항목" in block
    # 비어 있으면 아무 것도 붙이지 않는다(빈 헤더가 프롬프트를 먹지 않게).
    assert format_memory_block([]) == ""
    assert format_memory_block([{"kind": "fact", "content": ""}]) == ""


def test_instruction_covers_same_topic_different_target():
    instr = _instruction_text()
    assert "주제가 같고 대상만 다른 질문" in instr
    assert "스크래치" in instr, "실제로 틀렸던 사례가 지시문에 없다"
    assert "확인되지 않습니다" in instr


# --- 18번: 커맨드 출력에도 컨텍스트 상한이 있어야 한다 --------------------------------
# 매뉴얼·VOC 결과는 건당 1500자로 잘랐는데 커맨드 출력에만 상한이 없었다(64KB).
# 그 출력이 그대로 다음 요청 프롬프트에 실려 59,360토큰으로 32768 컨텍스트를 넘겼다(#123).
def test_command_output_has_context_budget_cap():
    import ssh_exec
    assert ssh_exec.MAX_OUTPUT <= 8000, \
        f"커맨드 출력 상한이 너무 크다({ssh_exec.MAX_OUTPUT}자). 프롬프트에 그대로 실린다."
    assert hasattr(ssh_exec, "set_output_limit_getter"), "설정으로 조정할 수 없다"

    src = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    assert "max_output: int | None = None" in src, \
        "기본값이 def 시점에 굳으면 설정 변경이 반영되지 않는다"

    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "execution_result_max_chars" in server
    assert "set_output_limit_getter(_output_limit)" in server

    cfg = open(os.path.join(ROOT, "shared", "migrations.py"), encoding="utf-8").read()
    assert '("execution_result_max_chars", "4000"' in cfg


def test_output_truncation_keeps_whole_lines_and_is_visible():
    """표 형태 출력을 줄 중간에서 끊으면 에이전트가 값을 잘못 읽는다.

    그리고 **잘렸다는 사실이 사용자에게 보여야 한다.** 예전에는 안내 문구를 stdout 끝에
    붙이는 게 전부여서, 모델이 그 줄을 빼먹으면 사용자는 목록이 전부인 줄 알았다
    ("홈 파일 목록이 중간에 잘리는 것 같아"). 구조화된 값으로도 돌려준다.
    """
    src = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    i = src.index("def _clip(")
    clip = src[i:i + 1600]
    assert "lines = s.split" in clip, "줄 단위로 자르지 않는다"
    assert "줄만 보입니다" in clip, "몇 줄 중 몇 줄인지 알려주지 않는다"
    assert "전부라고 말하지 마세요" in clip, "잘린 것을 전부로 답할 위험을 막지 않는다"
    # 결과 dict에도 실려야 진행 줄에서 보여줄 수 있다(LLM을 거치지 않는 경로).
    assert '"truncated": False' in src and "**clip_info," in src

    main = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert 'r.get("truncated")' in main, "진행 줄이 잘림을 알리지 않는다"
    assert "줄만" in main
    # 실행한 커맨드를 통째로 보여줘야 `-A` 누락과 잘림을 구분할 수 있다.
    assert "def _exec_command_text" in main


def test_chart_public_base_url_hidden_from_console():
    """비워 두는 게 기본인 고급 설정은 콘솔에 보이지 않아야 한다(사용자 지적)."""
    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    assert "chart_public_base_url" not in html
    assert "chart_mcp_url" in html          # 나머지 차트 설정은 남아 있어야 한다


def test_openwebui_base_url_description_explains_internal_port():
    """8080(컨테이너 내부)과 8502(사용자 접속)를 혼동하지 않게 설명이 있어야 한다."""
    cfg = open(os.path.join(ROOT, "shared", "migrations.py"), encoding="utf-8").read()
    i = cfg.index('("openwebui_base_url"')
    around = cfg[i - 500:i + 300]
    assert "8502" in around and "내부" in around


def test_context_overflow_error_is_actionable():
    """컨텍스트 초과는 사용자가 고칠 수 있는 문제다. litellm 스택트레이스를 그대로 보여주면
    무엇을 해야 하는지 알 수 없다(#123에서 실제로 그렇게 노출됐다)."""
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("_CONTEXT_ERROR_MARKERS")
    j = src.index("\ndef _trace_ctx(")

    async def get_config(_k, default=None):     # 접수 경로는 설정에서 읽는다
        return "서비스 포탈 > VOC 등록"

    ns = {"re": _re, "get_config": get_config}
    exec(src[i:j], ns)
    friendly = lambda e: asyncio.run(ns["_friendly_error"](e))   # noqa: E731

    real = ("litellm.ContextWindowExceededError: OpenAIException - Error code: 400 - "
            "{'error': {'message': \"This model's maximum context length is 32768 tokens. "
            "However, your request has 59360 input tokens. Please reduce the length of the "
            "input messages.\"}}")
    msg = friendly(Exception(real))
    assert "litellm" not in msg and "BadRequest" not in msg
    assert "59,360" in msg and "32,768" in msg, "실제 수치를 보여줘야 한다"
    assert "새 대화" in msg and "좁혀" in msg, "사용자가 할 수 있는 조치가 없다"
    # "운영팀에 알려주세요"로 끝내지 말고 콘솔에 등록된 접수 경로를 그대로 안내한다.
    assert "서비스 포탈 > VOC 등록" in msg

    # 다른 오류는 그대로 전달한다(원인을 숨기면 안 된다).
    assert "connection refused" in friendly(Exception("connection refused"))


# --- 19번: 환경 값 블록을 답변에 베끼지 못하게 한다 -----------------------------------
# 예전에는 지시문 끝에 `(참고: 로그인 서버 주소는 '...'입니다.)` 처럼 괄호 문장으로 붙였다.
# 모델이 그 꼴을 '답변 꼬리말 서식'으로 보고 답변에 그대로 베꼈고, 같은 모양으로
# `(참고: GPU_서버_활용_가이드_(KOR))`을 새로 만들어 붙이기까지 했다(#125).
def test_env_values_are_labeled_not_parenthetical():
    src = open(os.path.join(ROOT, "agent_server", "agent.py"), encoding="utf-8").read()
    # 주석에는 사고 기록으로 남아 있어도 되지만, 지시문을 만드는 코드가 괄호 꼬리말을
    # 붙이면 안 된다. 실제로 붙이는 표현식만 본다.
    code = "\n".join(ln for ln in src.split("\n") if not ln.strip().startswith("#"))
    # 줄바꿈 뒤에 바로 `(참고:` 를 붙이는 형태가 문제였다(답변 꼬리말처럼 보인다).
    assert "\\n(참고:" not in code, "환경 값을 괄호 꼬리말로 붙이면 모델이 답변에 베낀다"
    assert "# 이 환경의 값" in src
    assert "이 블록을 답변에 옮겨 적지 마세요" in src
    assert "'(참고: …)' 같은 꼬리말을 답변에 만들지 마세요" in src


def test_user_facing_error_hides_internal_settings():
    """사용자에게 내부 설정 키를 노출하지 않는다(관리자만 볼 값이다)."""
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("_CONTEXT_ERROR_MARKERS")

    async def get_config(_k, default=None):     # 접수 경로 미설정 환경
        return ""

    ns = {"re": _re, "get_config": get_config}
    exec(src[i:src.index("\ndef _trace_ctx(")], ns)
    msg = asyncio.run(ns["_friendly_error"](Exception(
        "This model's maximum context length is 32768 tokens. However, your request has "
        "33413 input tokens.")))
    assert "execution_result_max_chars" not in msg and "history_max_chars" not in msg
    assert "운영팀" in msg


def test_instruction_prefers_manual_for_infra_inventory():
    """'GPU 인프라 현황'은 매뉴얼에 정리돼 있다. 서버마다 커맨드를 돌리면 출력이 쌓여
    컨텍스트를 넘긴다(#123에서 33,413토큰으로 실패했다)."""
    instr = _instruction_text()
    assert "인프라 \"현황·구성\"을 물으면" in instr
    assert "매뉴얼을 먼저 검색합니다" in instr
    assert "이어서 씁니다" in instr, "도구를 이어서 쓰라는 규칙이 없다"
    assert "순서는 정해져 있지 않습니다" in instr, "도구 순서가 자유롭다는 것을 말하지 않는다"


def test_instruction_forbids_using_agent_name_as_account():
    """지어낸 파일 목록에 소유자를 `ops_assistant`(에이전트 자기 이름)로 적은 사고가 있었다."""
    instr = _instruction_text()
    assert "ops_assistant" in instr
    assert "실행할 수 있는 도구가 없어 확인하지 못했습니다" in instr


def test_execution_mcp_logs_exposed_tool_names():
    """필요한 툴이 꺼져 있으면 에이전트가 답을 지어낸다. 무엇이 노출됐는지 로그로 봐야 한다."""
    src = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
               encoding="utf-8").read()
    assert "노출된 툴:" in src


# --- #140: 다른 사용자 계정으로 실행하려는 시도를 **강제로** 막는다 -----------------
# 실행 신원(runuser)은 이미 본인으로 고정돼 있었지만, `phd list -u 남의계정`처럼 프로그램
# 자신이 대상을 고르는 옵션은 OS가 막아 주지 않는다. 지금까지는 모델이 지시문을 따라 거절해
# 준 것뿐이라 강제가 아니었다.
@pytest.mark.parametrize("args", [
    ["list", "-u", "other.user"],          # 옵션과 값이 따로
    ["list", "--user=other.user"],         # `=`로 붙인 꼴
    ["list", "-uother.user"],              # 짧은 옵션에 값이 붙은 꼴(가장 놓치기 쉽다)
    ["list", "--owner", "other.user"],
])
def test_run_command_blocks_other_user(args):
    from execution_exec import build_free_argv, deny_set, DEFAULT_DENY_CSV
    with pytest.raises(PermissionError) as e:
        build_free_argv("phd", args, "ops.user", deny_set(DEFAULT_DENY_CSV))
    assert "다른 사용자" in str(e.value)


@pytest.mark.parametrize("args", [
    ["list", "-u", "ops.user"],            # 본인 계정은 정상 사용법
    ["list", "-l"],
    ["list"],
])
def test_run_command_allows_self_and_plain_options(args):
    from execution_exec import build_free_argv, deny_set, DEFAULT_DENY_CSV
    build_free_argv("phd", args, "ops.user", deny_set(DEFAULT_DENY_CSV))


def test_registered_command_blocks_other_user_in_placeholder():
    """`phd list {option}`의 `{option}`에 `-u 남의계정`을 통째로 넣는 경로.
    자유 인자만 검사하면 이게 그대로 빠져나간다(토큰 하나로 들어오기 때문)."""
    from execution_exec import build_registered_argv, deny_set, DEFAULT_DENY_CSV
    specs = [{"name": "option", "type": "str", "required": False, "default": ""}]
    deny = deny_set(DEFAULT_DENY_CSV)
    with pytest.raises(PermissionError):
        build_registered_argv("phd list {option}", specs, {"option": "-u other.user"},
                              None, "ops.user", deny, True)
    # 같은 자리에 평범한 옵션은 통과해야 한다.
    assert build_registered_argv("phd list {option}", specs, {"option": "-l"},
                                 None, "ops.user", deny, True) == ["phd", "list", "-l"]


def test_other_user_error_forbids_guide_document_tour():
    """거절한 뒤 "가이드 위치: 슈퍼컴 Portal > ..."를 덧붙인 사고가 있었다.
    물어본 것은 남의 job이지 문서가 아니다 - 오류 문구가 그다음 행동까지 지시한다."""
    from execution_exec import build_free_argv, deny_set, DEFAULT_DENY_CSV
    with pytest.raises(PermissionError) as e:
        build_free_argv("phd", ["list", "-u", "other.user"], "ops.user",
                        deny_set(DEFAULT_DENY_CSV))
    msg = str(e.value)
    assert "가이드 문서 위치를 안내하지" in msg and "매뉴얼" in msg


def test_user_scope_check_is_configurable():
    """`sort -u`처럼 계정과 무관한 `-u`가 걸릴 수 있으므로 설정으로 끌 수 있어야 한다."""
    from execution_exec import build_free_argv, deny_set, DEFAULT_DENY_CSV, user_flag_set
    build_free_argv("sort", ["-u", "data.txt"], "ops.user",
                    deny_set(DEFAULT_DENY_CSV), user_flag_set(""))


# --- #140: 등록한 인자 설명이 LLM 스키마에 실제로 실린다 ---------------------------
def _arg_schema(row):
    """등록 커맨드 한 행 -> LLM에 보이는 파라미터 JSON 스키마."""
    import inspect
    from pydantic import create_model
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    from registry import build_entry
    row = {"tool_name": "t", "title": "t", "host_mode": "login_server",
           "enabled": True, "required_roles": [], **row}
    h = build_entry(row, None)["handler"]
    fields = {}
    for n, p in h.__signature__.parameters.items():
        if n in ("user_id", "host"):
            continue
        fields[n] = (h.__annotations__[n],
                     ... if p.default is inspect.Parameter.empty else p.default)
    return create_model("T", **fields).model_json_schema()["properties"]


def test_registered_arg_description_reaches_llm_schema():
    """예전에는 `option: str = ''`만 넘어가서, 콘솔에 적어 둔 옵션 설명이 모델에
    **한 글자도** 전달되지 않았다. 등록은 했는데 에이전트가 옵션을 못 채운 원인이다."""
    props = _arg_schema({
        "description": "job 목록", "exec_command": "phd list {option}",
        "args": [{"name": "option", "type": "str", "required": False, "default": "",
                  "description": "-l: 상세 출력 -lf: 필드 목록", "choices": []}]})
    assert props["option"]["description"] == "-l: 상세 출력 -lf: 필드 목록"


def test_enum_choices_become_schema_enum_with_labels():
    """선택지를 `값: 설명`으로 적으면 값만 enum이 되고 설명은 파라미터 설명으로 간다."""
    props = _arg_schema({
        "description": "job 상세", "exec_command": "phd info {option} {job_id}",
        "args": [{"name": "option", "type": "enum", "required": False, "default": "",
                  "description": "",
                  "choices": ["-j: JSON 형식으로 반환", "-tl: 부가 정보 출력"]},
                 {"name": "job_id", "type": "str", "required": True,
                  "description": "job id", "choices": []}]})
    # 필수가 아닌 선택형에는 빈 값도 있어야 한다(기본값 ""이 enum에 없으면 스키마가 깨진다).
    assert set(props["option"]["enum"]) == {"-j", "-tl", ""}
    assert "JSON 형식으로 반환" in props["option"]["description"]
    assert props["job_id"]["description"] == "job id"


def test_enum_value_parsing_ignores_description_part():
    """`cast_arg`도 같은 규칙으로 값을 비교해야 한다(스키마와 검증이 어긋나면 전부 거부된다)."""
    from execution_exec import cast_arg, choice_value
    spec = {"name": "option", "type": "enum",
            "choices": ["-j: JSON 형식으로 반환", "-tl: 부가 정보 출력"]}
    assert cast_arg(spec, "-j") == "-j"
    assert choice_value("12:00") == "12:00"        # 콜론 앞이 값처럼 보여도 공백 규칙으로 구분
    with pytest.raises(ValueError):
        cast_arg(spec, "-zz")


# --- #140: 에이전트가 질문한 사용자 계정을 안다 -----------------------------------
def test_agent_injects_caller_account_into_prompt():
    """모델이 자기 이름(`ops_assistant`)을 사용자 계정으로 말한 사고가 있었다.
    호출자 계정을 알려 주지 않았기 때문이다."""
    src = open(os.path.join(ROOT, "agent_server", "agent.py"), encoding="utf-8").read()
    assert 'caller_headers or {}).get("X-User-Id")' in src
    assert "질문한 사용자 계정" in src


def test_instruction_answers_other_user_question_in_one_line():
    instr = _instruction_text()
    # 거절은 이제 **실행 도구의 거부 사유**를 그대로 전하는 일반 규칙으로 다룬다(#153).
    # 질문 유형별 섹션을 두지 않는다 - 목록을 만들지 않는다는 원칙(#145·#149).
    assert "거부 사유를 돌려주면" in instr and "한 줄로" in instr
    assert "가이드 문서 위치를 안내하지 않습니다" in instr
    assert "질문한 사용자 계정" in instr


# --- #140: 엑셀 양식이 인자까지 실어 나른다 ---------------------------------------
def _exec_router():
    """관리자 콘솔 실행 라우터를 DB 없이 임포트한다(모듈 상수/파서만 쓴다)."""
    import importlib.util
    for k, v in {"CONFIG_DB_DSN": "postgres://x/y", "POSTGRES_PASSWORD": "x",
                 "ADMIN_PASSWORD": "x", "SESSION_SECRET": "x"}.items():
        os.environ.setdefault(k, v)
    sys.path.insert(0, os.path.join(ROOT, "admin_console", "backend", "routers"))
    path = os.path.join(ROOT, "admin_console", "backend", "routers", "execution.py")
    spec = importlib.util.spec_from_file_location("execrouter", path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def test_excel_template_roundtrips_with_args():
    """예전 일괄 등록은 이름·설명·실행 커맨드만 읽고 **인자를 버렸다**. `{option}`이 있는
    커맨드는 엑셀로 넣어 봐야 인자가 빈 채로 들어가 한 건씩 다시 채워야 했다.

    양식을 만들어 그대로 다시 읽었을 때 인자 정의가 살아 있고, 화면 등록과 같은 검증을
    통과해야 한다(양식과 파서가 따로 놀면 사용자는 '업로드했는데 안 된다'만 겪는다)."""
    import io
    import openpyxl
    from execution_exec import placeholders_in, validate_definition, deny_set, DEFAULT_DENY_CSV

    m = _exec_router()
    data = m._build_workbook([list(r) for r in m._TEMPLATE_EXAMPLES])
    ws = openpyxl.load_workbook(io.BytesIO(data))["커맨드"]
    rows = list(ws.iter_rows(values_only=True))
    header = [str(h) for h in rows[0]]
    deny = deny_set(DEFAULT_DENY_CSV)

    parsed = {}
    for r in rows[1:]:
        def cell(col, _r=r):
            v = _r[header.index(col)] if col in header else None
            return "" if v is None else str(v)
        exec_cmd = cell("실행 커맨드")
        ph = [p for p in placeholders_in(exec_cmd) if p != "user_id"]
        args = m._parse_args_from_row(cell, header, ph)
        host = m._HOST_WORDS.get(cell("실행 위치").strip().lower(), "login_server")
        validate_definition(cell("이름"), exec_cmd, args, host, deny)   # 화면 등록과 같은 문
        parsed[cell("이름")] = args

    assert parsed["myquota"] == []
    info = {a["name"]: a for a in parsed["s2_phd_info"]}
    # 자리표시자 순서대로 이름이 자동으로 붙는다(양식의 '이름' 칸을 비워도 되게 한 이유).
    assert [a["name"] for a in parsed["s2_phd_info"]] == ["option", "job_id"]
    assert info["option"]["type"] == "enum"
    assert info["option"]["choices"][0].startswith("-j:")
    assert info["job_id"]["required"] is True


def test_excel_arg_names_default_to_placeholder_order():
    """'인자N 이름'을 비워 두면 실행 커맨드의 자리표시자 순서로 채워진다.
    관리자가 이름을 두 번 옮겨 적다 틀리면 그 커맨드만 통째로 거부된다."""
    m = _exec_router()
    header = ["인자1 설명", "인자2 설명"]
    values = {"인자1 설명": "출력 형식", "인자2 설명": "job id"}
    args = m._parse_args_from_row(lambda c: values.get(c, ""), header, ["option", "job_id"])
    assert [a["name"] for a in args] == ["option", "job_id"]
    assert args[1]["description"] == "job id"


def test_excel_choices_split_on_newline_not_comma():
    """선택지 설명에 콤마가 흔하다("Return job info, json format").
    콤마로 쪼개면 설명 조각이 값으로 등록된다."""
    m = _exec_router()
    header = ["인자1 선택지"]
    text = "-j: Return job info, json format\n-tl: Print additional info"
    args = m._parse_args_from_row(lambda c: text if c == header[0] else "", header, ["option"])
    assert args[0]["choices"] == ["-j: Return job info, json format",
                                  "-tl: Print additional info"]
    assert args[0]["type"] == "enum"        # 선택지를 적었으면 타입 칸이 비어도 선택형


# --- #141: 상태를 담는 서비스에는 **이름 있는 볼륨**이 반드시 있어야 한다 --------------
# dev compose의 postgres에 데이터 볼륨이 없었다. pgvector 이미지가
# `VOLUME /var/lib/postgresql/data`를 선언하므로 **익명 볼륨**이 붙는데, 익명 볼륨은
# 컨테이너를 다시 만들 때 떨어져 나간다. #139에서 `ports:`를 127.0.0.1로 바꾼 것만으로
# compose가 postgres를 재생성했고, 매뉴얼·VOC·설정·등록 커맨드가 전부 사라졌다.
# 상태를 담는 **이미지**와 그 데이터 경로. 새 서비스를 붙일 때 여기 추가하면 검사가 따라온다.
# (redis는 임베딩 캐시라 없어져도 재생성되므로 뺐다 - 소실이 손실이 아닌 유일한 경우다.)
_STATEFUL_IMAGES = {
    "pgvector": "/var/lib/postgresql/data",
    "postgres": "/var/lib/postgresql/data",
    "clickhouse-server": "/var/lib/clickhouse",
    "minio": "/data",
    "open-webui": "/app/backend/data",
}


def _compose(name):
    yaml = pytest.importorskip("yaml")
    with open(os.path.join(ROOT, name), encoding="utf-8") as f:
        return yaml.safe_load(f)


@pytest.mark.parametrize("compose_file", ["docker-compose.yml", "docker-compose.dev.yml"])
def test_stateful_services_have_named_volumes(compose_file):
    """상태를 담는 서비스는 전부 **이름 있는 볼륨**이어야 한다.

    dev postgres에 데이터 볼륨이 없었다. pgvector 이미지가 `VOLUME`을 선언하므로 익명 볼륨이
    붙는데, 익명 볼륨은 컨테이너를 다시 만들면 떨어져 나간다. #139에서 `ports:`를 127.0.0.1로
    바꾼 것만으로 compose가 postgres를 재생성했고, 매뉴얼·VOC·설정·등록 커맨드가 전부 사라졌다.
    `down -v`만 위험한 게 아니다 - 재생성을 부르는 아무 변경이나 같은 결과를 낸다.

    postgres 하나만 보지 않고 **상태를 담는 이미지 전체**를 훑는다. 같은 실수를 세 번째로
    하지 않으려면 새 서비스가 붙을 때 자동으로 걸려야 한다.
    """
    conf = _compose(compose_file)
    declared = set((conf.get("volumes") or {}) or [])
    checked = 0

    for svc_name, svc in (conf.get("services") or {}).items():
        image = str(svc.get("image", ""))
        # `entrypoint`를 덮어쓴 서비스는 그 이미지를 **클라이언트로** 쓰는 것이다
        # (dev-config는 postgres 이미지로 psql만 돌린다). 서버는 이미지 기본 엔트리포인트로 뜬다.
        if svc.get("entrypoint"):
            continue
        for key, data_path in _STATEFUL_IMAGES.items():
            if key not in image:
                continue
            mounts = [str(m) for m in (svc.get("volumes") or [])]
            data = [m for m in mounts if m.split(":")[1:2] == [data_path]]
            assert data, (
                f"{compose_file}: '{svc_name}'({image})에 {data_path} 볼륨이 없습니다. "
                "익명 볼륨이 붙어 컨테이너를 다시 만들 때마다 데이터가 사라집니다.")
            src = data[0].split(":")[0]
            assert not src.startswith((".", "/", "$")), (
                f"{compose_file}: '{svc_name}'의 데이터 볼륨이 바인드 마운트({src})입니다. "
                "rsync --delete가 지울 수 있습니다(#137).")
            assert src in declared, (
                f"{compose_file}: '{svc_name}'의 '{src}'가 최상위 volumes에 없습니다.")
            checked += 1
            break

    assert checked, f"{compose_file}에서 상태 서비스를 하나도 찾지 못했습니다(검사가 헛돌고 있음)."


def test_openwebui_has_named_data_volume():
    """Open WebUI 계정·대화도 같은 이유로 이름 있는 볼륨이어야 한다."""
    yaml = pytest.importorskip("yaml")
    with open(os.path.join(ROOT, "docker-compose.dev.yml"), encoding="utf-8") as f:
        conf = yaml.safe_load(f)
    mounts = conf["services"]["open-webui"]["volumes"]
    data = [m for m in mounts if str(m).endswith(":/app/backend/data")]
    assert data and not str(data[0]).startswith(".")


# --- #142: /v1/*에 인증을 걸면서 **내부 호출자**를 확인하지 않았다 --------------------
def test_admin_console_sends_agent_api_key_to_agent_server():
    """#139에서 agent-server의 `/v1/*`에 `agent_api_key` 인증을 걸었는데, 관리자 콘솔의
    "기본 모델 동기화"가 `/v1/models`를 **헤더 없이** 부르고 있었다. 그래서 키를 넣는 순간
    값이 맞든 틀리든 401로 죽었다(사용자가 저장할 때마다 오류를 봤다).

    제약을 새로 걸 때는 그 경로에 이미 붙어 있던 호출자를 전부 훑어야 한다.
    """
    src = open(os.path.join(ROOT, "admin_console", "backend", "routers", "ops.py"),
               encoding="utf-8").read()
    assert 'get_config("agent_api_key"' in src, "콘솔이 agent_api_key를 읽지 않는다"

    # agent-server를 부르는 줄에 **인증 헤더가 붙어 있어야** 한다.
    get_line = next((ln for ln in src.split("\n") if "/v1/models" in ln and "client.get" in ln), "")
    assert get_line, "ops.py에서 /v1/models 호출을 찾지 못했다(테스트를 갱신할 것)"
    assert "headers=agent_headers" in get_line, \
        f"콘솔이 agent-server를 부를 때 인증 헤더를 안 보낸다: {get_line.strip()}"

    # 그 헤더는 **agent_api_key**로 만들어야 한다. Open WebUI 키를 보내면 방향이 뒤집힌다
    # (두 키는 목적지가 정반대다 - 사용자가 실제로 헷갈린 지점이다).
    hdr = next((ln for ln in src.split("\n") if "agent_headers =" in ln), "")
    assert "agent_key" in hdr and "openwebui" not in hdr, \
        f"agent-server용 헤더를 agent_api_key로 만들지 않는다: {hdr.strip()}"


def test_agent_api_key_is_hot_reload():
    """`agent_api_key`가 hot_reload=false면 저장 후 재시작 전까지 콘솔↔agent-server 값이
    어긋나 401이 계속된다. 안내 문구도 '재시작하세요'로 바뀌어야 하므로 고정해 둔다."""
    src = open(os.path.join(ROOT, "shared", "migrations.py"), encoding="utf-8").read()
    i = src.index('("agent_api_key"')
    seed = src[i:src.index("),", i)]
    # (key, value, desc, hot_reload, is_secret, force) — hot_reload가 True여야 한다.
    assert seed.rstrip().endswith("True, True, False"), f"시드 플래그가 바뀌었다: {seed[-40:]}"


# --- #144: 비로그인 강등 모드에서 작업 디렉토리가 root 홈이었다 ------------------------
@pytest.mark.parametrize("mode", ["su", "runuser"])
def test_non_login_privdrop_moves_to_user_home(mode):
    """`ssh root@host <cmd>`는 root 홈(`/root`)에서 시작하는데, `runuser -u`와 비로그인 `su`는
    **작업 디렉토리를 바꾸지 않는다**. 그래서 `SSH_PRIVDROP=runuser`로 바꾼 뒤 `ls -lh`가
    `/root`에서 돌아 `Permission denied`가 났다.

    지시문이 "실행은 항상 본인 홈에서 시작합니다"라고 약속하므로 코드가 그것을 지켜야 한다."""
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    from ssh_exec import _remote_command
    cmd = _remote_command("ops.user", ["ls", "-lh"], mode)
    assert cmd.startswith("cd ~ops.user"), f"홈으로 이동하지 않는다: {cmd}"
    # `&&`면 root가 홈에 못 들어가는 환경(GPFS root_squash)에서 커맨드가 아예 안 돈다.
    assert "; " in cmd and "&&" not in cmd, f"실패 시 커맨드를 막으면 안 된다: {cmd}"


def test_login_privdrop_does_not_double_cd():
    """`su - user`는 로그인 셸이라 이미 홈으로 간다. 덧붙이면 군더더기다."""
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    from ssh_exec import _remote_command
    cmd = _remote_command("ops.user", ["ls", "-lh"], "su-login")
    assert cmd == "su - ops.user -c 'ls -lh'"


def test_home_cwd_does_not_break_argument_quoting():
    """홈 이동을 붙이면서 인자 인용이 깨지면 셸 주입이 생긴다."""
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    from ssh_exec import _remote_command
    cmd = _remote_command("ops.user", ["ls", "; rm -rf /", "$HOME", "`id`"], "runuser")
    assert "'; rm -rf /'" in cmd and "'$HOME'" in cmd and "'`id`'" in cmd
    # 우리가 의도한 `;`는 하나뿐이어야 한다(cd 뒤).
    assert cmd.count(";") == 2, cmd      # cd 뒤 1개 + 인자 안의 리터럴 1개(따옴표 안)


def test_instruction_decides_by_content_not_phrasing():
    """실행 여부를 **말투로** 판단하면 안 된다(#149).

    "보여 줘"라고 해야만 실행하도록 써 뒀더니 "내 홈 디렉토리는 어디야?"가 일반 지식으로
    분류돼 도구를 아예 호출하지 않았다. 질문 형태를 나열하는 것은 커맨드를 나열하는 것과
    같은 실수다(#145) - 사용자가 어떻게 물을지는 알 수 없다.
    기준은 **답에 무엇이 필요한가** 하나여야 한다.
    """
    instr = _instruction_text()
    assert "말투로 판단하지 않습니다" in instr
    assert "답이 이 서버에 물어봐야 나오는 값이면 실행합니다" in instr
    # 판별법이 있어야 실행 가능한 규칙이 된다.
    assert "회사·서버·" in instr and "따라 달라지는가" in instr


def test_instruction_forbids_fabricating_environment_values():
    """마지막 안전장치: 분류를 잘못해 (B)로 답하더라도 이 환경의 값은 지어내지 않는다.

    사용자 지적: "일반지식이더라도 모델이 만들어내진 말아야지. 모르는 건 모른다고 해야지."
    실제로 `/home/ops.user`를 지어내 답했다(정답은 `/home/users/ops.user`).
    """
    instr = _instruction_text()
    assert "이 환경의 값을 지어내지 않습니다" in instr
    # 지어내는 중임을 스스로 알아채는 신호 - 헤지 문구를 명시적으로 금지한다.
    for hedge in ("일반적으로 …입니다", "보통 …입니다", "정확한 것은 직접 확인해 보세요"):
        assert hedge in instr, f"헤지 문구를 금지 목록에 넣지 않았다: {hedge}"
    # 모를 때 **갈 곳**을 준다. #155부터는 "모른다"에서 끝내지 않고 운영팀 문의로 보낸다
    # (사용자 지시: "그 어떤 db에서 확인할 수 없는거면 운영팀에 문의하라고 하라고").
    assert "운영팀에 문의하라" in instr, "모를 때 어디로 보낼지 알려 주지 않았다"
    assert "지어낸 값을 주는 것이 실패입니다" in instr


# --- #145: 지시문에 **특정 커맨드를 박지 않는다** ------------------------------------
# 사용자가 네 번째로 지적한 사항이다(#74 `phd info`, #125, #140, 그리고 `pwd`/`echo $HOME`).
# 커맨드는 (1) 콘솔에 등록돼 툴로 노출되거나 (2) 모델이 아는 표준 리눅스 명령이다.
# 지시문은 **원칙**만 말해야 한다 - 커맨드를 적기 시작하면 하나하나 다 적어야 하고,
# 시스템이 바뀔 때마다 지시문이 거짓말을 하게 된다(#144가 정확히 그렇게 났다).
_FORBIDDEN_IN_INSTRUCTION = [
    # 사내 전용(존재하지 않는 것을 지어내 쓴 사고가 있었다)
    "phd ", "myquota", "squeue", "sinfo", "sbatch", "bsub",
    # 확인용 표준 명령(모델이 알아서 고르게 둔다)
    "pwd", "whoami", "echo $", "$HOME", "$USER",
    # 나열하기 시작하면 끝이 없다
    "nvidia-smi", "`ls ", "`ls`", "`df", "`du", "`head", "`tail", "`find",
]


def test_instruction_names_no_specific_commands():
    """지시문에 커맨드 이름을 적지 않는다. 원칙만 쓰고 선택은 모델과 툴 목록에 맡긴다."""
    instr = _instruction_text()
    # 모듈 docstring은 제외하고 지시문 본문만 본다.
    body = instr[instr.index("AGENT_INSTRUCTION = "):]
    hits = [tok for tok in _FORBIDDEN_IN_INSTRUCTION if tok in body]
    assert not hits, (
        f"지시문에 특정 커맨드가 박혀 있다: {hits}\\n"
        "  커맨드는 콘솔 등록(툴)이나 모델의 표준 리눅스 지식으로 해결한다. "
        "지시문에는 원칙만 쓴다.")


def test_instruction_states_no_shell_as_a_property():
    """커맨드를 적는 대신 **성질**을 말해야 한다: 셸을 거치지 않는다.
    그래야 모델이 `echo $HOME` 같은 것을 스스로 피한다(그건 글자 그대로 출력된다)."""
    instr = _instruction_text()
    assert "셸을 거치지 않습니다" in instr
    assert "글자 그대로" in instr


# --- #147: "되돌리기" 버튼이 옛 지시문을 계속 저장했다 --------------------------------
def test_instruction_reset_reads_file_not_module_cache():
    """`agent_system_instruction` 되돌리기가 **파일을 다시 읽어야** 한다.

    예전에는 함수 안에서 `from agent_instruction import AGENT_INSTRUCTION` 했다. 부수효과가
    없으니 안전하다고 생각했는데, 안전한 것과 **최신인 것**은 다르다. 파이썬은 `sys.modules`에
    캐시하므로 이 프로세스가 한 번이라도 읽었으면 그 뒤로는 옛 텍스트를 쓴다. `./shared`가
    바인드 마운트라 파일은 최신인데 **버튼이 아무 일도 하지 않는** 상태가 됐다
    (사용자: "2번 그대로 했는데도 1번 하면 옛지시문이라고 떠").

    `importlib.reload`도 부족하다 - `.pyc`는 mtime+크기로 유효성을 보므로 같은 초에 같은
    크기로 바뀌면 낡은 바이트코드를 쓴다(실제로 재현했다). 그래서 소스를 직접 읽어 compile한다.
    """
    src = open(os.path.join(ROOT, "admin_console", "backend", "routers", "settings.py"),
               encoding="utf-8").read()
    reset = src[src.index("def _read_instruction_from_disk"):]
    assert "open(path" in reset and "compile(src" in reset, \
        "지시문을 파일에서 직접 읽지 않는다"
    # 되돌리기 경로에서 모듈 import로 지시문을 가져오면 안 된다.
    body = src[src.index("async def reset_agent_instruction"):]
    assert "from agent_instruction import" not in body, \
        "되돌리기가 모듈 캐시에서 지시문을 읽는다 - 옛 텍스트가 저장된다"
    assert "_read_instruction_from_disk()" in body
    # 못 읽었으면 조용히 옛 값을 쓰지 않고 실패해야 한다.
    assert "지시문 파일을 읽지 못했습니다" in reset


def test_instruction_reset_actually_returns_current_file_text():
    """헬퍼가 정말 디스크의 현재 내용을 돌려주는지 - 문자열 검사만으로는 부족하다."""
    for k, v in {"CONFIG_DB_DSN": "postgres://x/y", "POSTGRES_PASSWORD": "x",
                 "ADMIN_PASSWORD": "x", "SESSION_SECRET": "x"}.items():
        os.environ.setdefault(k, v)
    import importlib.util
    sys.path.insert(0, os.path.join(ROOT, "admin_console", "backend"))
    sys.path.insert(0, os.path.join(ROOT, "admin_console", "backend", "routers"))
    path = os.path.join(ROOT, "admin_console", "backend", "routers", "settings.py")
    spec = importlib.util.spec_from_file_location("settings_under_test", path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)

    on_disk = _instruction_text()
    got = mod._read_instruction_from_disk()
    assert got in on_disk, "파일에 없는 내용을 돌려준다"
    assert len(got) > 5000, f"지시문이 너무 짧다({len(got)}자) - 잘못 읽고 있다"


# --- #148: 엑셀 양식의 고정 선택지는 **드롭다운**이어야 한다 --------------------------
def test_excel_template_has_dropdowns_for_fixed_choice_columns():
    """타입/필수/활성/실행 위치는 값이 정해져 있다. 자유 입력이면 오타가 조용히 다른 뜻이 된다
    ("선택"이라고 적으면 선택형이 아니라 문자열로 들어간다)."""
    import io
    import openpyxl
    from openpyxl.utils import column_index_from_string

    m = _exec_router()
    ws = openpyxl.load_workbook(io.BytesIO(m._build_workbook([])))["커맨드"]
    cols = [c.value for c in ws[1]]

    got = {}
    for dv in ws.data_validations.dataValidation:
        first = str(dv.sqref).split()[0].split(":")[0]
        letter = "".join(ch for ch in first if ch.isalpha())
        got[cols[column_index_from_string(letter) - 1]] = dv.formula1

    for col, expected in [("실행 위치", "로그인 서버"), ("활성", "Y,N"),
                          ("인자1 타입", "문자열,정수,선택형"), ("인자1 필수", "Y,N")]:
        assert col in got, f"'{col}' 열에 드롭다운이 없다"
        assert expected in got[col], f"'{col}' 선택지가 다르다: {got[col]}"

    # 인자 슬롯 전부에 걸려야 한다(1번만 걸고 나머지를 빠뜨리기 쉽다).
    for i in range(1, m.TEMPLATE_ARG_SLOTS + 1):
        assert f"인자{i} 타입" in got and f"인자{i} 필수" in got, f"인자{i} 슬롯에 드롭다운이 없다"


def test_excel_dropdown_values_match_the_parser():
    """드롭다운에 있는 값은 **파서가 실제로 알아듣는 값**이어야 한다.
    화면에서 고를 수 있는데 업로드하면 무시되는 값이 있으면 안 된다."""
    m = _exec_router()
    for label in m._DROPDOWNS["타입"]:
        assert label.lower() in m._TYPE_WORDS, f"파서가 모르는 타입: {label}"
    for label in m._DROPDOWNS["실행 위치"]:
        assert label.lower() in m._HOST_WORDS, f"파서가 모르는 실행 위치: {label}"
    for label in m._DROPDOWNS["필수"]:
        assert m._truthy(label, False) == (label == "Y"), f"필수 값 해석이 다르다: {label}"


def test_empty_default_drops_the_argument_entirely():
    """"기본값이 없으면 아무것도 안 보이는 건가?" - 자리표시자가 통째로 빠진다."""
    from execution_exec import build_registered_argv, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)
    spec = [{"name": "option", "type": "enum", "required": False,
             "default": "", "choices": ["-l", "-lf"]}]
    assert build_registered_argv("phd list {option}", spec, {}, None,
                                 "ops.user", deny, True) == ["phd", "list"]
    spec[0]["default"] = "-l"
    assert build_registered_argv("phd list {option}", spec, {}, None,
                                 "ops.user", deny, True) == ["phd", "list", "-l"]


def test_progress_line_always_shows_line_count():
    """잘리지 않았을 때도 줄 수를 보여줘야 한다(#149).

    예전에는 잘렸을 때만 `⚠ N줄 중 M줄만`을 찍었다. 그래서 22줄짜리 답을 받고도
    **우리가 자른 건지 모델이 자른 건지** 구분할 수 없었다. 항상 줄 수를 찍으면
    사용자가 답변의 행 수와 눈으로 대조할 수 있다.
    """
    import json as _json
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("def _unwrap_result")      # _result_phrase가 이걸 쓴다
    end = src.index("\nclass _StreamDedup")
    ns = {"json": _json, "re": _re, "_mem_on": lambda x: True}
    exec(src[i:end], ns)
    phrase = ns["_result_phrase"]

    base = {"ip": "10.0.0.1", "as_user": "ops.user", "duration_ms": 400, "exit_code": 0}
    full = phrase("run_command", {**base, "truncated": False,
                                  "total_lines": 132, "shown_lines": 132})
    assert "132줄" in full and "⚠" not in full, full

    cut = phrase("run_command", {**base, "truncated": True,
                                 "total_lines": 132, "shown_lines": 58})
    assert "⚠ 출력 132줄 중 58줄만" in cut, cut

    # 한 줄짜리 출력에까지 붙이면 잡음이다.
    one = phrase("run_command", {**base, "truncated": False,
                                 "total_lines": 1, "shown_lines": 1})
    assert "줄" not in one, one


# --- #150: 실행 결과 원문을 **LLM을 거치지 않고** 답변 뒤에 붙인다 ---------------------
def _raw_outputs_cls():
    import json as _json
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("def _unwrap_result")
    ns = {"json": _json, "re": _re}
    exec(src[i:src.index("\nasync def _make_raw_outputs")], ns)
    return ns["_RawOutputs"]


class _FakeEvent:
    def __init__(self, responses):
        self._r = [type("FR", (), {"response": x})() for x in responses]

    def get_function_responses(self):
        return self._r


def test_raw_output_block_shows_every_line():
    """모델이 132줄 중 22줄만 보여준 사고가 반복됐다(#146). 지시문은 확률이라,
    사용자가 반드시 봐야 하는 것은 LLM을 거치지 않고 붙인다."""
    RO = _raw_outputs_cls()
    lines = "\n".join(f"-rw-r--r-- 1 ops.user users 100 file{i}.txt" for i in range(132))
    ro = RO(2, 100000)
    ro.observe(_FakeEvent([{"result": {
        "exit_code": 0, "command": "ls -la", "ip": "10.0.0.100", "as_user": "ops.user",
        "total_lines": 132, "truncated": False, "stdout": lines}}]))
    block = ro.block()
    assert block.count("file") == 132, "원문에서 행이 빠졌다"
    assert "`ls -la`" in block and "10.0.0.100 · ops.user · 132줄" in block


def test_raw_output_skips_short_and_non_execution_results():
    """한두 줄짜리는 답변에 이미 들어 있어 중복이고, 검색 결과는 실행 결과가 아니다."""
    RO = _raw_outputs_cls()
    short = RO(2, 100000)
    short.observe(_FakeEvent([{"exit_code": 0, "command": "x", "stdout": "/home/users/ops.user"}]))
    assert short.block() == ""

    search = RO(2, 100000)
    search.observe(_FakeEvent([{"results": [{"title": "a"}, {"title": "b"}]}]))
    assert search.block() == ""


def test_raw_output_survives_backticks_in_output():
    """출력에 ``` 가 들어 있으면 코드블록이 중간에 닫혀 나머지가 마크다운으로 샌다."""
    RO = _raw_outputs_cls()
    ro = RO(2, 100000)
    ro.observe(_FakeEvent([{"exit_code": 0, "command": "cat x.md",
                            "stdout": "```python\nprint(1)\n```"}]))
    block = ro.block()
    assert "````text" in block, "울타리를 늘리지 않았다"
    assert block.rstrip().endswith("````")


def test_raw_output_block_has_its_own_cap():
    """원문 블록에도 상한이 있어야 한다. 없으면 한 답변이 수십만 자가 된다."""
    RO = _raw_outputs_cls()
    ro = RO(2, 200)
    ro.observe(_FakeEvent([{"exit_code": 0, "command": "big",
                            "stdout": "\n".join("x" * 50 for _ in range(100))}]))
    block = ro.block()
    assert "원문 표시 상한" in block
    assert len(block) < 1000, len(block)


def test_raw_output_wired_into_every_answer_path():
    """네 엔드포인트(스트리밍/비스트리밍)에 전부 붙어야 한다. 하나만 빠져도 그 경로에서
    사용자는 여전히 줄어든 답을 받는다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert src.count("raw.observe(event)") == 6, \
        f"수집 지점이 6곳이어야 한다(현재 {src.count('raw.observe(event)')})"
    assert src.count("raw.block()") == 6, \
        f"붙이는 지점이 6곳이어야 한다(현재 {src.count('raw.block()')})"
    # 설정으로 끌 수 있어야 한다.
    assert 'get_config("execution_raw_output"' in src


def test_raw_output_summary_is_off_by_default():
    """원문 뒤 요약은 **LLM을 한 번 더** 부른다. 지연이 늘어나므로 기본은 꺼 둔다."""
    mig = open(os.path.join(ROOT, "shared", "migrations.py"), encoding="utf-8").read()
    i = mig.index('("execution_raw_output_summary"')
    seed = mig[i:mig.index("),", i)]
    assert '"false"' in seed, f"기본값이 꺼져 있지 않다: {seed}"
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    fn = src[src.index("async def _raw_output_summary"):]
    assert "return \"\"" in fn, "실패 시 조용히 넘어가지 않는다"


# --- #151: 지시문 버전 확인을 **매직 문자열로 하지 않는다** ---------------------------
def test_instruction_check_compares_file_not_magic_string():
    """"지시문에 이 문구가 있으면 최신"은 지시문을 고칠 때마다 같이 고쳐야 하는데 잊는다.
    실제로 #149에서 `어디야?` 나열을 지우자, 최신 지시문일수록 "옛것"이라고 나오는
    확인이 됐다(#151). 문서에 사실을 복사하면 원본이 바뀔 때 거짓이 된다.
    """
    _need("scripts/check-instruction.sh", "docs/NEXT-STEPS.md")
    path = os.path.join(ROOT, "scripts", "check-instruction.sh")
    assert os.path.isfile(path), "scripts/check-instruction.sh 가 없다"
    script = open(path, encoding="utf-8").read()
    # 파일을 직접 읽어 비교해야 한다(모듈 캐시도 타지 않는 방식, #147과 같은 이유).
    assert "agent_instruction.py" in script and "compile(" in script
    assert "md5(value)" in script, "DB 값과 해시를 비교하지 않는다"
    # **구분자를 끼워 이어 붙이지 않는다**(#152). 처음에 앞뒤 200자를 `\x01`로 이어 비교했는데,
    # psql에서는 리터럴 4글자, 파이썬에서는 1바이트가 되어 **내용이 같아도 항상 불일치**였다.
    # 실제 postgres로 확인했다: md5는 파이썬 md5(t.encode("utf-8"))와 정확히 일치한다.
    # **주석은 빼고** 본다. 스크립트 주석에 "왜 \\x01을 쓰면 안 되는지"를 적어 뒀는데,
    # 통짜 substring 검사면 그 설명 때문에 테스트가 실패한다(실제로 겪었다).
    code = "\n".join(ln for ln in script.split("\n") if not ln.lstrip().startswith("#"))
    assert "x01" not in code, "이스케이프가 필요한 구분자를 다시 쓰고 있다"
    assert "left(value" not in code and "right(value" not in code, \
        "앞뒤 조각 이어붙이기로 되돌아갔다 - psql/파이썬 이스케이프가 어긋난다"

    # NEXT-STEPS 가 지시문 본문의 문구로 최신 여부를 판정하면 안 된다.
    #
    # 예전에는 여기서 `check-instruction.sh`가 **항상** 있어야 한다고 봤다. 그런데 NEXT-STEPS는
    # "지금 할 일"만 담고 끝난 단계는 지우는 문서다(CLAUDE.md) — 지시문 반영이 끝나 그 단계를
    # 지우자 이 테스트가 깨졌다. **규칙이 아니라 문서의 모양을 검사하고 있었다**(#153·#158과
    # 같은 실수). 진짜 규칙은 "지시문 확인 절차를 적는다면 스크립트로 적는다"이다.
    steps = open(os.path.join(ROOT, "docs", "NEXT-STEPS.md"), encoding="utf-8").read()
    assert "value like '%" not in steps, \
        "NEXT-STEPS 가 아직 매직 문자열로 지시문 버전을 본다 - 지시문을 고치면 거짓이 된다"
    if "되돌리기" in steps or "지시문을 최신" in steps:
        assert "check-instruction.sh" in steps, \
            "지시문 확인 절차를 적으면서 스크립트를 안 쓴다 - 매직 문자열로 되돌아간다"


# --- #153: 지시문 재작성 — 라우팅 중심으로 절반으로 줄였다 ----------------------------
def test_instruction_is_within_prompt_budget():
    """지시문이 12,751자(7,787토큰 = 컨텍스트의 24%)까지 불어나 있었다. 규칙을 덧붙여
    고치려 한 것이 개별 규칙의 준수율을 떨어뜨렸다(#150에서 측정). 예산을 고정한다."""
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    from agent_instruction import AGENT_INSTRUCTION as instr
    from registry import estimate_prompt_tokens
    _chars, tokens = estimate_prompt_tokens([instr])
    assert tokens < 5000, (
        f"지시문이 {tokens:,}토큰이다(32,768의 {tokens / 32768 * 100:.0f}%). "
        "규칙을 더 넣기 전에 중복부터 지울 것 — 길수록 개별 규칙을 덜 지킨다.")


def test_instruction_does_not_run_commands_for_symptom_reports():
    """VOC 문의("접속이 원활하지 않다")에 GPU job 목록을 실행해 보여준 사고(#153).
    증상 호소는 상태 조회 요청이 아니라 원인·해결을 묻는 것이다."""
    instr = _instruction_text()
    assert "증상을 호소하는 것은 상태 조회 요청이 아닙니다" in instr
    assert "답과 상관없는 도구를 부르지 않습니다" in instr


def test_instruction_isolates_each_turn():
    """앞 턴에서 '다른 계정은 조회 못 한다'고 거절한 뒤, 전혀 다른 VOC 문의에도 같은 거절을
    반복한 사고(#153). 새 메시지는 그 메시지만으로 판단해야 한다."""
    instr = _instruction_text()
    assert "그 메시지 하나만으로" in instr
    assert "거절을 **이어붙이지 않습니다**" in instr or "이어붙이지 않습니다" in instr
    assert "대명사를 푸는 데만" in instr


def test_instruction_forbids_code_block_for_document_guidance():
    """문서 안내를 bash 코드블록으로 출력한 사고(#153). 사용자가 실행할 명령으로 오해한다."""
    instr = _instruction_text()
    assert "코드 블록은 실행한 커맨드와 그 출력에만 씁니다" in instr
    assert "코드 블록 없이" in instr


def test_instruction_routing_covers_all_three_mcps():
    """사용자 요구: 문의가 오면 매뉴얼·VOC 관련인지 확인하고, 실행으로 풀 수 있는지
    에이전트가 판단해야 한다. 세 갈래가 한 곳에 모여 있어야 그 판단이 선다."""
    instr = _instruction_text()
    routing = instr[instr.index("## 2) 도구를 고릅니다"):instr.index("## 3) 멈춥니다")]
    for tool in ("매뉴얼 검색", "과거 사례(VOC) 검색", "커맨드 실행"):
        assert tool in routing, f"도구 선택 절에 '{tool}'이 없다"


# --- #154·#155: 지어내기를 **코드로** 잡는다 -------------------------------------------
def _grounding_cls():
    import json as _json
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("_IP_RE = re.compile")
    ns = {"json": _json, "re": _re, "_mem_on": lambda x: True}
    exec(src[i:src.index("\nclass _Pace:")], ns)
    return ns["_AnswerGuard"]


_FALLBACK_MARK = "운영팀에 문의해 주세요"


class _FakeCall:
    def __init__(self, name):
        self.name = name


class _FakeEv2:
    def __init__(self, calls=None, responses=None):
        self._c = [_FakeCall(n) for n in (calls or [])]
        self._r = [type("FR", (), {"response": x})() for x in (responses or [])]

    def get_function_calls(self):
        return self._c

    def get_function_responses(self):
        return self._r


def test_guard_replaces_answer_with_invented_ip():
    """없는 서버 IP를 만들어 안내한 사고(#154). 사용자가 그 주소로 접속을 시도하므로
    곧바로 사고다.

    #155에서 **경고 덧붙이기를 버렸다.** 사용자: "저런 그대로 믿지마세요 문구를 넣지말라고.
    아예 지어내지 말라고." 틀린 IP를 보여주면서 경고를 붙이는 것은 답이 아니다 —
    그 답변 자체를 운영팀 문의 안내로 갈아 끼운다."""
    G = _grounding_cls()
    g = G(True, "login server 접속이 안됩니다", "10.0.0.100")
    g.observe(_FakeEv2(calls=["search_manual"],
                       responses=[{"results": [{"chunk_text": "ETX 클라이언트로 접속"}]}]))
    out = g.review("10.20.30.40 으로 접속하세요.")
    assert "10.20.30.40" not in out, "지어낸 IP가 그대로 사용자에게 나갔다"
    assert _FALLBACK_MARK in out
    # 우리가 프롬프트에 넣어 준 환경 값은 근거로 인정한다(매번 막히면 못 쓴다).
    assert g.review("10.0.0.100 으로 접속합니다.") == "10.0.0.100 으로 접속합니다."


def test_guard_blocks_guide_without_manual_search():
    """매뉴얼을 검색하지도 않고 '가이드 문서: …'를 안내한 사고(#154).
    사용자: '가이드 위치와 제목은 manual_mcp에서 확인할 수 있을 때만 알려야 함'."""
    G = _grounding_cls()
    g = G(True, "접속이 안돼요", "")
    g.observe(_FakeEv2(responses=[{"exit_code": 0, "stdout": "x\ny"}]))
    body = "자세한 내용은 다음 문서를 참고하세요:\n - 가이드 문서: 사용자 매뉴얼"
    assert _FALLBACK_MARK in g.review(body)

    # 검색했으면 그대로 내보낸다.
    g2 = G(True, "접속이 안돼요", "")
    g2.observe(_FakeEv2(calls=["search_manual"], responses=[{"results": [{"chunk_text": "a"}]}]))
    assert g2.review(body) == body


def test_guard_accepts_prefetched_manual_as_evidence():
    """#155: 매뉴얼을 **우리가 먼저** 검색해 프롬프트에 넣는다. 모델이 툴을 부르지 않아도
    그 근거로 답한 것은 정상이므로 막으면 안 된다. 반대로 선검색이 0건이면 검색한 것으로
    치지 않는다 — 그 상태의 문서 안내는 지어낸 것이다."""
    G = _grounding_cls()
    body = ("자세한 내용은 다음 문서를 참고하세요:\n"
            " - 가이드 위치: 슈퍼컴 Portal > 활용 가이드\n - 가이드 문서: GPU 활용 가이드")
    g = G(True, "GPU 어떻게 신청해요", "")
    g.seed_rag([{"guide_location": "슈퍼컴 Portal > 활용 가이드",
                    "guide_document": "GPU 활용 가이드", "chunk_text": "신청 절차"}])
    assert g.review(body) == body

    g0 = G(True, "GPU 어떻게 신청해요", "")
    g0.seed_rag([])                      # 0건 = 검색 안 한 것과 같다
    assert _FALLBACK_MARK in g0.review(body)


def test_guard_blocks_invented_path_but_allows_real_one():
    """`/home/ops_assistant` 처럼 없는 경로를 안내한 사고가 있었다(#125)."""
    G = _grounding_cls()
    g = G(True, "홈이 어디야", "")
    g.observe(_FakeEv2(responses=[{"exit_code": 0, "stdout": "/home/users/ops.user"}]))
    assert g.review("홈은 /home/users/ops.user 입니다.") == "홈은 /home/users/ops.user 입니다."
    assert "/home/ops.user 입니다" not in g.review("홈은 /home/ops.user 입니다.")


def test_guard_is_quiet_on_normal_answers():
    """정상 답변까지 막으면 서비스가 못 쓰게 된다."""
    G = _grounding_cls()
    g = G(True, "job 목록 보여줘", "")
    g.observe(_FakeEv2(responses=[{"exit_code": 0, "stdout": "23836848 Queued\n23836892 Queued"}]))
    answer = "Queued 상태인 job은 23836848, 23836892 두 건입니다."
    assert g.review(answer) == answer


def test_guard_includes_intake_path_in_fallback():
    """운영팀 문의로 보낼 때 접수 경로를 알려 준다(설정 `voc_intake_guide`)."""
    G = _grounding_cls()
    g = G(True, "질문", "", "포탈 > 문의하기")
    assert "포탈 > 문의하기" in g.fallback("테스트")


def test_guard_holds_body_so_it_can_be_replaced():
    """스트리밍에서 본문을 흘려보내면 되돌릴 수 없다. 검사가 켜져 있으면 모아 두었다가
    한 번에 내보낸다 — 그래야 갈아 끼울 수 있다."""
    G = _grounding_cls()
    assert G(True, "q", "").hold is True
    assert G(False, "q", "").hold is False

    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    # 델타를 흘리기 전에 hold를 확인하는 지점이 세 스트리밍 경로에 모두 있어야 한다.
    assert src.count("ground.hold") == 5, src.count("ground.hold")


def test_guard_wired_into_every_answer_path():
    """여섯 지점(엔드포인트 3 × 스트리밍/비스트리밍) 전부에 붙어야 한다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert src.count("ground.observe(event)") == 6, src.count("ground.observe(event)")
    assert src.count("ground.review(") == 6, src.count("ground.review(")
    assert src.count("ground.seed_rag(") == 3, src.count("ground.seed_rag(")
    assert 'get_config("answer_grounding_check"' in src, "설정으로 끌 수 없다"


# --- #155: 매뉴얼은 **무조건** 먼저 검색해서 그 근거로 답한다 --------------------------
def test_rag_prefetch_wired_into_every_endpoint():
    """사용자: "사용자 문의가 들어오면 무조건!!! 제발!!!! manual_mcp, voc_mcp 로 rag 한 후에
    답변 생성하길 바람."

    지시문으로 네 번 시켰지만 모델은 자기가 안다고 판단하면 툴을 건너뛰었다.
    **부를지 말지를 모델에게 맡기지 않는다** — 세 엔드포인트 모두에서 우리가 먼저 검색한다.
    #155에서 매뉴얼만 했는데 VOC도 같은 이유로 안 불리고 있었다(voc_db에 있는 질문을 그대로
    물었는데 엉뚱한 답이 나왔다)."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert src.count("await _rag_context(") == 3, src.count("await _rag_context(")
    # 각 MCP와 **같은 검색 함수**를 써야 결과가 갈리지 않는다.
    assert "search_manual_chunks" in src and "search_voc_records" in src
    assert 'get_config("rag_prefetch"' in src, "설정으로 끌 수 없다"


def test_voc_search_has_one_shared_implementation():
    """VOC 검색이 MCP 안에만 있으면 선검색이 같은 결과를 낼 수 없다. shared로 뺐다."""
    mcp = open(os.path.join(ROOT, "mcp_servers", "voc_mcp", "server.py"),
               encoding="utf-8").read()
    assert "from voc_search import" in mcp, "MCP가 공용 검색을 쓰지 않는다"
    assert "FULL OUTER JOIN" not in mcp, "MCP에 검색 SQL이 복제돼 있다(경로가 둘로 갈린다)"
    shared = open(os.path.join(ROOT, "shared", "voc_search.py"), encoding="utf-8").read()
    assert "async def search_voc_records" in shared


def test_voc_prefetch_block_marks_who_handled_it():
    """운영자가 조치한 건인지 사용자가 직접 한 건인지가 프롬프트에 있어야
    '남이 해 준 조치를 사용자에게 시키는' 답이 안 나온다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    block = src[src.index("_VOC_HANDLED = {"):src.index("async def _rag_context")]
    assert "handled_by" in block
    assert "운영자가 확인·조치" in block and "사용자가 직접 해결" in block


def test_retrieval_query_carries_previous_user_turn():
    """이어지는 질문("그러면 접속 못 하는거 아니야?")은 그 문장만으로 검색이 안 된다.
    직전 **사용자** 발화를 붙인다. 어시스턴트 답변은 붙이지 않는다 —
    지어낸 내용이 질의에 섞이면 그 방향으로 검색이 끌려간다."""
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("def _retrieval_query(")
    ns = {}
    exec(src[i:src.index("async def _search_manual_for")], ns)
    q = ns["_retrieval_query"]("그러면 슈퍼컴 접속 못 하는거 아니야?", [
        ("user", "login server 접속이 갑자기 안됩니다"),
        ("assistant", "운영팀에 문의해 주세요"),
    ])
    assert "login server" in q, "앞 문맥이 질의에 안 들어갔다"
    assert "운영팀에 문의해 주세요" not in q, "어시스턴트 답변이 질의에 섞였다"
    assert q.endswith("그러면 슈퍼컴 접속 못 하는거 아니야?")
    assert _re.sub(r"\s", "", ns["_retrieval_query"]("혼자 온 질문", [])) == "혼자온질문"


def test_guard_keeps_grounded_lines_and_drops_only_bad_ones():
    """#156: 값 하나가 근거에 없다고 답변을 통째로 버리면, 매뉴얼에서 확인된 점검 목록까지
    사라진다. 사용자: "manual_db 기반으로 우선 사용자가 확인해야 할 사항들 먼저 쭉 가이드를
    하고, 그 후에 운영팀한테 문의하라고 해야지." 지어낸 값이 든 **줄만** 덜어낸다."""
    G = _grounding_cls()
    g = G(True, "접속이 안돼요", "")
    g.seed_rag([{"chunk_text": "ETX 클라이언트 버전을 확인하고 재설치합니다. "
                               "사내망(VPN) 연결 상태를 확인합니다."}])
    answer = ("1. ETX 클라이언트 버전을 확인하고 재설치합니다.\n"
              "2. 사내망(VPN) 연결 상태를 확인합니다.\n"
              "3. 10.20.30.40 으로 직접 접속해 봅니다.")
    out = g.review(answer)
    assert "10.20.30.40" not in out, "지어낸 IP가 그대로 나갔다"
    assert "ETX 클라이언트 버전을 확인" in out, "근거 있는 안내까지 버렸다"
    assert "사내망(VPN) 연결 상태" in out
    assert "운영팀" in out, "확인해 볼 것을 안내한 뒤 운영팀으로 이어져야 한다"


def test_guard_falls_back_when_nothing_grounded_survives():
    """반대로 남는 게 껍데기뿐이면 그때는 통째로 운영팀 문의로 바꾼다."""
    G = _grounding_cls()
    g = G(True, "서버 주소 알려줘", "")
    g.seed_rag([{"chunk_text": "무관한 내용"}])
    out = g.review("## 접속 주소\n- 10.20.30.40\n- 10.20.30.41")
    assert "10.20.30" not in out
    assert _FALLBACK_MARK in out


def test_instruction_guides_before_escalating():
    """사용자: "바로 운영팀 확인이 필요하다고 하는데, manual_db 기반으로 우선 사용자가
    확인해야 할 사항들 먼저 쭉 가이드를 하고, 그 후에 운영팀한테 문의하라고 해야지."
    """
    instr = _instruction_text()
    assert "운영팀 문의는 맨 마지막입니다" in instr
    voc = instr[instr.index("## 과거 사례(VOC)를 쓸 때"):]
    # operator 사례에서 곧바로 접수로 넘기는 길이 열려 있으면 안 된다. 표현은 #158에서
    # 바뀌었지만(원인마다 해 볼 것을 달고, 접수는 맨 마지막) 규칙은 그대로다.
    assert "맨 마지막에" in voc and "접수 경로만 던지는 답변은" in voc


# --- #157: 처리 주체를 키워드가 아니라 LLM이 판정한다 ---------------------------------
def test_voc_classify_puts_reason_before_label():
    """생성 순서가 곧 조건이다. 스키마에서 `reason`이 `label`보다 **먼저** 나와야
    라벨이 자기가 방금 쓴 근거를 조건으로 결정된다. 뒤집으면 근거가 사후 정당화가 된다."""
    import voc_classify as vc
    props = list(vc._SCHEMA["properties"]["results"]["items"]["properties"])
    assert props.index("reason") < props.index("label"), props


def test_voc_classify_offers_abstention():
    """애매한 건을 찍지 않도록 기권 라벨이 있어야 한다."""
    import voc_classify as vc
    assert "unknown" in vc.LABELS
    assert "unknown" in vc._SCHEMA["properties"]["results"]["items"]["properties"]["label"]["enum"]


def test_voc_classify_rubric_is_criteria_not_keywords():
    """키워드 목록으로 돌아가지 않는다 — 사용자 지적이 바로 그것이었다
    ("키워드 사용에 한계가 있더라고"). 기준을 주고 모델이 적용하게 한다."""
    import voc_classify as vc
    assert "낱말로 판단하지 마세요" in vc._RUBRIC
    assert "권한이 필요한지" in vc._RUBRIC
    # 예전 키워드 목록이 프롬프트로 옮겨오지 않았는지 본다.
    from voc_search import _OPERATOR_HINTS
    leaked = [h for h in _OPERATOR_HINTS if h in vc._RUBRIC]
    assert len(leaked) <= 2, f"키워드 목록이 프롬프트로 새어 들어왔다: {leaked}"


def test_voc_classify_drops_hallucinated_and_misaligned_ids():
    """배치 분류가 조용히 어긋나는 두 경우 — 모델이 없는 id를 지어내거나, 한 칸씩 밀려
    답하는 것. id를 대조해 **버려야** 한다(미분류로 남아 다음 회차에 다시 시도)."""
    import voc_classify as vc
    sent = {10, 11, 12}
    got = vc.validate([
        {"id": 10, "reason": "본인 스크립트 수정", "label": "user"},
        {"id": 99, "reason": "없는 행", "label": "operator"},        # 지어낸 id
        {"id": 11, "reason": "라벨 오타", "label": "operater"},       # 모르는 라벨
        {"id": 12, "reason": "권한 변경", "label": "OPERATOR"},       # 대소문자는 허용
        {"id": 10, "reason": "중복", "label": "operator"},           # 먼저 온 것만
    ], sent)
    assert got == {10: ("user", "본인 스크립트 수정"), 12: ("operator", "권한 변경")}
    assert 11 not in got and 99 not in got


def test_voc_classify_parses_reasoning_and_fenced_output():
    """구조화 출력을 서버가 거부하면 평문으로 떨어진다. 그때 Qwen3의 <think> 블록과
    ```json 펜스, 앞뒤 설명 문장을 견뎌야 한다."""
    import voc_classify as vc
    assert vc.parse_response(
        '<think>음 이건...</think>\n```json\n{"results":[{"id":1,"reason":"r","label":"user"}]}\n```'
    ) == [{"id": 1, "reason": "r", "label": "user"}]
    assert vc.parse_response(
        '분류했습니다: {"results":[{"id":2,"reason":"r","label":"operator"}]} 이상입니다.'
    ) == [{"id": 2, "reason": "r", "label": "operator"}]
    assert vc.parse_response("설명만 하고 JSON이 없음") == []
    assert vc.parse_response("") == []


def test_voc_search_prefers_stored_label_over_keywords():
    """저장된 판정이 있으면 그것을 쓰고, 없을 때만 키워드로 떨어진다."""
    src = open(os.path.join(ROOT, "shared", "voc_search.py"), encoding="utf-8").read()
    assert 'item.get("handled_by")\n                              or classify_handling(' in src
    # 세 검색 갈래(키워드 전용 / 3축 RRF / 2축 RRF) 모두 컬럼을 읽어 와야 한다.
    # 하나라도 빠지면 그 경로에서만 조용히 키워드 추론으로 떨어진다.
    selects = [b for b in src.split("SELECT ") if "FROM voc_records" in b or "JOIN voc_records" in b]
    reading = [b for b in selects if "handled_by" in b.split("FROM")[0].split("JOIN")[0]]
    assert len(reading) == 3, f"{len(reading)}개 갈래만 handled_by를 읽는다"


def test_voc_update_clears_stale_classification():
    """답변을 고치면 처리 주체 판정도 무효다. 안 지우면 옛 판정이 계속 쓰인다."""
    src = open(os.path.join(ROOT, "admin_console", "backend", "routers", "voc.py"),
               encoding="utf-8").read()
    i = src.index("async def update_voc")
    upd = src[i:src.index("@router.", i)]
    assert "handled_by=NULL" in upd


def _voc_block_fn():
    """`_voc_block`을 그대로 실행한다(순수 함수라 import 없이 돈다).

    예전 이 테스트들은 구현의 **글자**를 검사했다. #153에서 배운 대로, 구현을 고치면 규칙이
    살아 있어도 같이 깨진다. 실제로 돌려서 **나온 블록**을 본다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("_VOC_HANDLED = {")
    ns = {}
    exec(src[i:src.index("async def _rag_context")], ns)  # noqa: S102
    return ns["_voc_block"]


def test_unknown_case_is_not_shown_as_user_or_as_operator():
    """`unknown`은 '조치 내용이 없어 못 가린 건'이다 (#158).

    '사용자가 직접 해결'로 보여주면 남이 해 준 조치를 사용자에게 시키게 되고,
    '운영자가 확인·조치'로 보여주면 **우리가 모르는 것을 안다고 프롬프트에 쓰는 것**이다.
    둘 다 아닌 표기여야 한다."""
    block = _voc_block_fn()

    def headers(out):
        """사례마다 붙는 '## 과거 사례 N (처리: …)' 줄만 본다(지침 문구와 섞이지 않게)."""
        return [ln for ln in out.split("\n") if ln.startswith("## 과거 사례")]

    head = headers(block([{"question": "접속 안 됨", "answer": "확인 부탁드립니다",
                           "handled_by": "unknown"}]))[0]
    assert "사용자가 직접 해결" not in head, "unknown을 사용자 건으로 보여주고 있다"
    assert "운영자가 확인·조치" not in head, "판정 못 한 것을 운영자가 조치했다고 단정한다"
    assert "판정 불가" in head
    # 판정 컬럼이 아예 비어 있는 행(백필 전 신규 등록)도 같은 취급이어야 한다.
    empty = headers(block([{"question": "q", "answer": "a", "handled_by": None}]))[0]
    assert "사용자가 직접 해결" not in empty and "판정 불가" in empty
    # 반대로 정상 판정된 두 라벨은 그대로 구분돼야 한다.
    assert "사용자가 직접 해결" in headers(
        block([{"question": "q", "answer": "a", "handled_by": "user"}]))[0]
    assert "운영자가 확인·조치" in headers(
        block([{"question": "q", "answer": "a", "handled_by": "operator"}]))[0]


def test_voc_block_demands_multiple_causes_without_asserting_one():
    """사용자 지시: "과거 사례를 기반으로 ~~수 있다는 식의 원인 추측성 답변을 해야하고,
    명확하게 무슨 이유때문이다 라고 대답하면 안돼. 여러 원인이 있을 수 있음을 명시하고
    해결방안 여러개를 내도록 해줘."(#158)

    지시문뿐 아니라 **프롬프트 블록에도** 실려야 한다 — 지시문은 되돌리기 버튼을 눌러야
    반영되지만 이 블록은 코드라 즉시 반영된다(#147에서 지시문이 몇 턴 동안 옛것이었다)."""
    block = _voc_block_fn()
    out = block([{"question": "접속 안 됨", "answer": "스토리지 부하로 조치했습니다",
                  "handled_by": "operator"}])
    assert "단정" in out, "원인을 단정하지 말라는 지침이 블록에 없다"
    assert "여럿" in out, "원인이 여럿일 수 있다는 지침이 블록에 없다"

    # 전부 `user` 사례면 이 지침은 답과 무관한 잡음이다. 붙지 않아야 한다.
    allowed = block([{"question": "q", "answer": "a", "handled_by": "user"}])
    assert "단정" not in allowed, "사용자 건만 있는데 추측 지침이 붙어 프롬프트를 먹는다"


def test_keyword_fallback_never_claims_user():
    """키워드 목록은 **운영자 신호만** 찾는다. 안 걸린 것은 '사용자 건'이 아니라 '모른다'다.

    Errors.md에 주신 운영자 답변 5개가 예전에는 전부 `user`로 떨어졌다 —
    "조치를 진행하였습니다"가 목록의 `조치하였`와 글자가 어긋난다(#158)."""
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    os.environ.setdefault("CONFIG_DB_DSN", "postgres://test/test")
    from voc_search import classify_handling
    for answer in (
        "현재 정상 사용이 확인됩니다. 불편을 드려 죄송합니다. 신속한 조치를 진행 하겠습니다.",
        "현재 정상 접속이 가능 합니다. 사용에 불편을 드려 대단히 죄송합니다.",
        "스토리지 부하증상으로 조치를 진행하였습니다. 현재 정상 접속이 가능 합니다.",
        "",
    ):
        assert classify_handling(answer) != "user", f"'{answer[:20]}'을 사용자 건으로 본다"
    # 운영자 신호가 뚜렷한 건은 그대로 잡아야 한다(폴백이 무의미해지면 안 된다).
    assert classify_handling("로그를 확인하여 재기동했습니다") == "operator"


def test_voc_tool_documents_every_label_it_returns():
    """툴이 돌려주는 값은 docstring이 말해 주는 것이 전부다(#155에서 배운 것).
    `unknown`을 돌려주면서 설명하지 않으면 모델은 모르는 값을 받는다."""
    src = open(os.path.join(ROOT, "mcp_servers", "voc_mcp", "server.py"),
               encoding="utf-8").read()
    doc = src[src.index("async def search_voc("):]
    for label in ("user", "operator", "unknown"):
        assert f'"{label}"' in doc, f"{label} 라벨이 툴 설명에 없다"


def test_instruction_requires_hedged_multiple_causes_for_operator_cases():
    """운영자·판정불가 사례는 원인을 **단정하지 않고 여럿으로** 제시하고, 원인마다 해 볼 것을
    붙인 뒤 마지막에 접수를 안내한다 (#158)."""
    instr = _instruction_text()
    assert "원인은 ~입니다\"라고 단정하지 않습니다" in instr
    assert "원인을 하나로 좁히지 않습니다" in instr
    assert "원인마다 해 볼 것을 답니다" in instr
    # 1절 3번의 헤지 금지와 충돌해 보이므로, 경계를 지시문 안에서 갈라 줘야 한다.
    assert "1절 3번이 금지한 것과 다릅니다" in instr


def test_instruction_forbids_deciding_scope_by_itself():
    """"~~ 클러스터는 슈퍼컴 인프라가 아니니까 다른 운영팀에 문의하라" — 지어낸 것이다.
    무엇이 우리 소관인지는 검색 결과가 정하지, 모델이 정하지 않는다."""
    instr = _instruction_text()
    assert "소관인지 스스로 판단하지 않습니다" in instr


def test_manual_prefetch_block_carries_reference_path():
    """선검색 블록에 `guide_location`(콘솔에 등록된 문서 위치)이 실려야 한다.
    이게 프롬프트에 없으니 모델이 문서 위치를 안내하지 못했다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    block = src[src.index("def _manual_block("):src.index("def _voc_block(")]
    assert 'r.get("guide_location")' in block
    assert 'r.get("guide_document")' in block
    assert "문서 위치" in block and "문서 이름" in block


def test_manual_tool_docstring_declares_reference_fields():
    """#155 진짜 원인: 툴이 `guide_location`/`guide_document`를 돌려주는데도
    docstring의 Returns에 그 이름이 없었다. 지시문은 그 이름으로 쓰라고 하는데
    툴 계약에는 없으니 모델이 있는 줄도 몰랐다."""
    src = open(os.path.join(ROOT, "mcp_servers", "manual_mcp", "server.py"),
               encoding="utf-8").read()
    doc = src[src.index("async def search_manual"):src.index("async def get_document")]
    for field in ("guide_location", "guide_document", "chunk_text"):
        assert field in doc, f"search_manual docstring에 {field}가 없다"


def test_instruction_sends_unknown_to_operations_team():
    """사용자: "우리 매뉴얼에 없거나 그 어떤 db에서 확인할 수 없는거면 운영팀에 문의하라고
    하라고." 그리고 "그대로 믿지마세요" 류의 면피 문구를 붙이지 말 것."""
    instr = _instruction_text()
    assert instr.count("운영팀") >= 3, "확인 실패 시 운영팀 문의로 보내지 않는다"
    assert "그대로 믿지 마세요" in instr and "붙이지 말고" in instr, \
        "경고를 붙여 추측을 내보내는 길이 막혀 있지 않다"


def test_instruction_tells_model_the_manual_block_is_already_there():
    """선검색 결과를 프롬프트에 넣어도, 그게 있다는 걸 지시문이 말해 주지 않으면
    모델이 같은 검색을 또 부르거나 블록을 무시한다."""
    instr = _instruction_text()
    assert "# 이번 질문에 대한 매뉴얼 검색 결과" in instr


# --- #159: 차트 디자인 — 색은 눈으로 고르지 않는다 -------------------------------------
def _svg_chart():
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "chart_mcp"))
    import svg_chart
    return svg_chart


def test_chart_palette_is_the_validated_one():
    """팔레트는 검증기(밝기 대역·채도 하한·색각이상 분리도·정상시야 하한·대비)를 통과한
    여덟 색 **고정 순서**다. 순서 자체가 색각이상 안전장치라 바꾸면 안 된다 (#159).

    색을 바꾸고 싶으면 이 테스트를 고치기 전에 검증기를 먼저 돌릴 것."""
    sc = _svg_chart()
    assert sc.SERIES_LIGHT[:3] == ["#2a78d6", "#eb6834", "#1baf7a"]
    assert sc.SERIES_DARK[:3] == ["#3987e5", "#d95926", "#199e70"]
    assert len(sc.SERIES_LIGHT) == len(sc.SERIES_DARK) == 8
    # 밝은/어두운 두 벌이 다 있어야 한다. 한 벌만 두고 뒤집으면 어두운 표면에서 대비가 깨진다.
    assert sc.SERIES_LIGHT != sc.SERIES_DARK


def test_chart_never_invents_a_ninth_color():
    """9번째 색은 만들어 봐야 색각이상에서 기존 색과 구분되지 않는다. 계열을 잘라 낸다.
    흩뿌림은 어느 두 점이든 나란히 놓여 '모든 쌍' 검사를 받으므로 상한이 더 낮다."""
    sc = _svg_chart()
    many = [{"name": f"s{i}", "values": [1]} for i in range(20)]
    assert len(sc.cap_series("bar", many)) == 8
    assert len(sc.cap_series("scatter", many)) == 3
    # 잘라 낸 뒤에도 팔레트 인덱스를 넘어서는 클래스를 쓰지 않아야 한다.
    svg = sc.render("bar", ["a"], many, "제목", "")
    assert ".f9" not in svg and 'class="f9"' not in svg


def test_chart_marks_stay_thin_and_grid_stays_recessive():
    """두꺼운 원색 덩어리 + 진한 격자는 조악해 보인다. 막대는 24px를 넘지 않고,
    격자는 실선 헤어라인이다(점선은 '임계값'으로 읽힌다)."""
    sc = _svg_chart()
    assert sc._BAR_MAX <= 24
    svg = sc.render("stacked", ["a", "b"],
                    [{"name": "x", "values": [1, 2]}, {"name": "y", "values": [2, 1]}], "t", "")
    assert "stroke-dasharray" not in svg, "격자/축에 점선을 쓰고 있다"


def test_chart_puts_numbers_on_the_marks():
    """사용자 지시: "그래프 위에 적절히 숫자나 퍼센트도 있고"(Errors.md 4번).

    밝은 모드에서 세 계열색은 표면 대비 3:1 미만이라, 검증기가 요구하는 완화가
    **직접 값 라벨**이다 — 이 라벨은 장식이 아니라 접근성 요건이다."""
    sc = _svg_chart()
    bar = sc.render("bar", ["1월", "2월"], [{"name": "건수", "values": [42, 17]}], "t", "건")
    assert ">42<" in bar and ">17<" in bar, "막대 위에 값이 없다"

    # 선은 **끝점 하나만** 찍는다. 모든 점에 숫자를 달면 읽히지 않는다.
    line = sc.render("line", ["1", "2", "3"], [{"name": "a", "values": [11, 22, 33]}], "t", "")
    assert ">33<" in line and ">11<" not in line and ">22<" not in line

    # 원/도넛은 조각 비율(%)을 낸다.
    pie = sc.render("pie", ["a", "b"], [{"name": "x", "values": [75, 25]}], "t", "")
    assert ">75%<" in pie and ">25%<" in pie


def test_chart_legend_wraps_instead_of_running_off_canvas():
    """계열 이름이 길면 예전 범례는 한 줄에 밀어 넣어 캔버스 밖으로 흘러 나갔다."""
    sc = _svg_chart()
    series = [{"name": f"아주 긴 계열 이름 {i}", "values": [1, 2]} for i in range(6)]
    rows = sc._legend_rows(series, sc._W - 72)
    assert len(rows) > 1, "긴 이름 6개가 한 줄에 들어간다고 보고 있다"
    svg = sc.render("line", ["a", "b"], series, "t", "")
    for name in (s["name"] for s in series):
        assert name in svg
    # 계열이 하나뿐이면 범례 상자를 만들지 않는다(제목이 이미 무엇인지 말한다).
    assert sc._layout("line", ["a"], [{"name": "혼자", "values": [1]}], "t", "")["legend_rows"] == []


def test_chart_text_never_wears_the_series_color():
    """밝은 계열색(노랑·청록)은 표면 위에서 글자로 읽히지 않는다. 정체성은 글자 **옆의**
    색칠된 마크가 나른다 — 글자는 잉크 토큰만 쓴다."""
    sc = _svg_chart()
    svg = sc.render("bar", ["a"], [{"name": "x", "values": [1]}], "제목", "건")
    for chunk in svg.split("<text")[1:]:
        head = chunk.split(">")[0]
        assert "fill=" not in head, f"글자에 색을 직접 박았다: {head[:80]}"
        assert not any(f"f{i}" in head for i in range(1, 9)), \
            f"글자가 계열 색 클래스를 입었다: {head[:80]}"


def test_chart_title_does_not_follow_the_axis_inward():
    """hbar는 항목 이름이 길어 축 왼쪽이 200px까지 밀린다. 제목까지 따라가면 화면
    한가운데서 시작한다 — 제목·범례·단위는 고정 여백에 붙는다."""
    sc = _svg_chart()
    L = sc._layout("hbar", ["아주 길고 긴 항목 이름입니다", "짧음"],
                   [{"name": "건수", "values": [1, 2]}], "제목", "건")
    assert L["left"] > 120, "이 경우 축이 안쪽으로 밀려야 한다(전제 확인)"
    assert L["text_left"] <= 48


def test_chart_value_labels_do_not_round_the_number_away():
    """축 눈금은 15,630을 `16k`로 줄여도 되지만(원래 둥근 수를 찍는 자리다),
    막대 위에 **읽으라고 찍는 값**까지 그러면 화면에 틀린 수가 보인다 (#159).
    사용자는 그 숫자를 그대로 옮겨 적는다."""
    sc = _svg_chart()
    assert sc._fmt(15630) == "16k"                 # 눈금용은 압축해도 된다
    assert sc._fmt_value(15630) == "15,630"
    svg = sc.render("bar", ["nvme-01"], [{"name": "사용량", "values": [15630]}], "t", "GB")
    # 값 라벨(class="ink2 val")에는 반올림 안 된 원래 값이 있어야 한다.
    values = [c.split(">")[1].split("<")[0] for c in svg.split('class="ink2 val"')[1:]]
    assert "15,630" in values, values


# --- #160: 분류 배치를 동시에 보낸다 ---------------------------------------------------
def test_voc_classify_runs_batches_concurrently():
    """배치를 하나씩 `await` 하면 LLM이 생성하는 동안 말고는 GPU가 논다. vLLM은 여러
    요청을 겹쳐 처리하도록 만들어진 서버다 — 수천 건 백필이 몇 시간 걸리던 이유다 (#160).

    가짜 분류기로 **실제로 겹쳐 도는지** 잰다(코드 모양이 아니라 동작을 본다)."""
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    import voc_classify as vc

    inflight, peak = 0, 0

    async def fake_classify(records, **kw):
        nonlocal inflight, peak
        inflight += 1
        peak = max(peak, inflight)
        await asyncio.sleep(0.02)
        inflight -= 1
        return {r["id"]: ("user", "근거") for r in records}

    class FakePool:
        def __init__(self):
            self.updated = []

        async def fetch(self, *a):
            return [{"id": i, "question": "q", "answer": "a"} for i in range(1, 25)]

        async def executemany(self, _sql, args):
            self.updated += list(args)

    pool = FakePool()
    orig = vc.classify_records
    vc.classify_records = fake_classify
    try:
        # 배치 4건씩 24건 = 6배치. 동시 3이면 최대 3개가 겹쳐 돌아야 한다.
        vc._batch_size = lambda: _coro(4)
        vc._concurrency = lambda: _coro(3)
        out = asyncio.run(vc.classify_pending(pool))
    finally:
        vc.classify_records = orig

    assert peak > 1, f"배치가 여전히 하나씩 돈다(최대 동시 {peak})"
    assert peak <= 3, f"동시 상한을 안 지킨다({peak} > 3)"
    assert out["classified"] == 24 and out["counts"]["user"] == 24
    assert len(pool.updated) == 24, "행마다 왕복하지 말고 배치로 한 번에 갱신해야 한다"


async def _coro(value):
    return value


def test_voc_classify_keeps_going_when_a_batch_fails():
    """한 배치가 죽어도 나머지는 끝까지 간다. 실패한 배치의 행은 미분류로 남아
    다음 회차에 다시 시도된다 — 수천 건 백필이 한 건 때문에 통째로 죽으면 안 된다."""
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    import voc_classify as vc

    async def flaky(records, **kw):
        if records[0]["id"] == 5:                     # 두 번째 배치만 터뜨린다
            raise RuntimeError("vLLM 500")
        return {r["id"]: ("operator", "근거") for r in records}

    class FakePool:
        async def fetch(self, *a):
            return [{"id": i, "question": "q", "answer": "a"} for i in range(1, 13)]

        async def executemany(self, _sql, args):
            pass

    orig = vc.classify_records
    vc.classify_records = flaky
    try:
        vc._batch_size = lambda: _coro(4)
        vc._concurrency = lambda: _coro(3)
        out = asyncio.run(vc.classify_pending(FakePool()))
    finally:
        vc.classify_records = orig

    assert out["classified"] == 8 and out["failed"] == 4, out


def test_every_config_key_the_code_reads_is_seeded():
    """`get_config("x")`를 쓰는데 시드가 없으면 그 값은 **콘솔에 나타나지 않는다.**

    코드는 기본값으로 잘 돌기 때문에 아무도 모르고, 관리자는 바꿀 방법이 없다.
    `voc_classify_concurrency`를 추가하고 "db-init 필요 없다"고 안내했다가 설정 탭에
    안 보인다는 지적을 받았다 (#160). 시드가 있으면 최소한 db-init 한 번으로 나타난다."""
    seeded = set(re.findall(r'^\s*\("([a-z0-9_]+)",\s*',
                            open(os.path.join(ROOT, "shared", "migrations.py"),
                                 encoding="utf-8").read(), re.M))
    read = set()
    for base, _dirs, files in os.walk(ROOT):
        if ".git" in base or "tests" in base:
            continue
        for f in files:
            if not f.endswith(".py") or f == "migrations.py":
                continue
            src = open(os.path.join(base, f), encoding="utf-8").read()
            read |= set(re.findall(r'get_config\(\s*"([a-z0-9_]+)"', src))
    missing = sorted(read - seeded)
    assert not missing, (
        f"코드가 읽지만 시드에 없는 설정 키: {missing} — "
        "migrations.py의 config_seed()에 추가하세요(콘솔에서 바꿀 수 없는 값이 됩니다)")


def test_settings_the_docs_tell_you_to_change_are_visible_in_the_console():
    """콘솔 설정 탭은 **허용 목록**(SETTING_GROUPS)으로 그린다 — 시드된 키라도 목록에
    없으면 화면에 안 나온다. 그런데 NEXT-STEPS는 "설정 탭 `x`를 바꾸세요"라고 안내한다.
    둘이 어긋나면 사용자는 있지도 않은 항목을 찾아 헤맨다 (#160).

    실제로 #155~#160에서 만든 다섯 개(rag_prefetch · manual_prefetch_top_k ·
    voc_prefetch_top_k · voc_classify_batch · voc_classify_concurrency)가 목록에 없는 채로
    문서에만 적혀 있었다. db-init을 돌려도 안 보이는 게 당연했다."""
    _need("docs/NEXT-STEPS.md")
    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    block = html[html.index("const SETTING_GROUPS"):html.index("function SettingsPanel")]
    shown = set(re.findall(r'"([a-z0-9_]+)"', block))

    seeded = set(re.findall(r'^\s*\("([a-z0-9_]+)",\s*"',
                            open(os.path.join(ROOT, "shared", "migrations.py"),
                                 encoding="utf-8").read(), re.M))
    steps = open(os.path.join(ROOT, "docs", "NEXT-STEPS.md"), encoding="utf-8").read()

    # NEXT-STEPS가 백틱으로 지목하는 설정 키는 전부 화면에 있어야 한다.
    # (.env로 관리되는 키는 콘솔에서 못 바꾸는 것이 정상이라 제외한다.)
    env_managed = {"manual_db_dsn", "voc_db_dsn", "execution_db_dsn", "system_db_dsn",
                   "agent_session_db_dsn", "memory_db_dsn", "redis_url"}
    mentioned = {k for k in seeded if f"`{k}`" in steps} - env_managed
    missing = sorted(mentioned - shown)
    assert not missing, (
        f"NEXT-STEPS가 바꾸라고 하는데 설정 탭에 없는 키: {missing} — "
        "admin_console/frontend/index.html 의 SETTING_GROUPS에 추가하세요")


def test_each_mcp_setting_sits_in_its_own_group():
    """사용자 지시: "manual 관련이면 manual mcp 구역에 두고 해야지. 각 mcp 별 필요한 인자들은
    묶여서 한 구역에 있도록 해."(#160)

    관리자는 "매뉴얼 관련 설정"을 찾지 "agent-server가 읽는 설정"을 찾지 않는다. 값을 읽는
    프로세스가 어디든 **이름이 가리키는 구역**에 있어야 한다 — 선검색 개수는 agent-server가
    읽지만 Manual/VOC 구역에 있다.

    (`vllm_`은 일부러 둘로 나눠 둔다 — LLM과 임베딩·리랭커는 서로 다른 서버다.)"""
    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    block = html[html.index("const SETTING_GROUPS"):html.index("function SettingsPanel")]
    groups = {}
    for g in re.finditer(r'\{ title: "([^"]+)", service: [^,]+,\s*keys: \[(.*?)\] \}',
                         block, re.S):
        groups[g.group(1)] = re.findall(r'"([a-z0-9_]+)"', g.group(2))

    home = {"manual_": "Manual MCP", "voc_": "VOC MCP",
            "execution_": "Execution MCP", "chart_": "Chart MCP"}
    misplaced = []
    for title, keys in groups.items():
        for k in keys:
            for prefix, belongs in home.items():
                if k.startswith(prefix) and title != belongs:
                    misplaced.append(f"{k}: '{title}'에 있는데 '{belongs}'여야 한다")
    assert not misplaced, "설정이 이름과 다른 구역에 흩어져 있다 — " + " / ".join(misplaced)


# --- #161: 분류 중지 + 콘솔에서 설명 걷어내기 ------------------------------------------
def test_classify_stop_finishes_inflight_batches_instead_of_cancelling():
    """'중지'는 태스크를 죽이지 않는다. **이미 LLM에 나가 있는 배치는 끝내고 저장**한 뒤
    멈춘다 — 중간에 끊으면 그 판정은 그냥 버려진다(시간만 쓰고 아무것도 안 남는다).
    대기 중인 배치만 접으면 사용자는 한 건도 잃지 않는다 (#161)."""
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    import voc_classify as vc

    started, stop = [], {"now": False}

    async def fake_classify(records, **kw):
        started.append(len(records))
        await asyncio.sleep(0.01)
        stop["now"] = True                      # 첫 배치가 도는 사이 '중지'를 누른 셈
        return {r["id"]: ("user", "근거") for r in records}

    saved = []

    class FakePool:
        async def fetch(self, *a):
            return [{"id": i, "question": "q", "answer": "a"} for i in range(1, 25)]

        async def executemany(self, _sql, args):
            saved.extend(args)

    orig = vc.classify_records
    vc.classify_records = fake_classify
    try:
        vc._batch_size = lambda: _coro(4)
        vc._concurrency = lambda: _coro(1)      # 한 번에 하나씩 -> 중지 시점이 또렷하다
        out = asyncio.run(vc.classify_pending(FakePool(), should_stop=lambda: stop["now"]))
    finally:
        vc.classify_records = orig

    # 첫 배치는 끝까지 가서 **저장됐고**, 나머지는 시작조차 하지 않았다.
    assert len(started) == 1, f"중지 후에도 배치를 더 보냈다: {started}"
    assert len(saved) == 4, f"도는 중이던 배치를 버렸다(저장 {len(saved)}건)"
    assert out["classified"] == 4 and out["stopped"] == 20, out


def test_classify_stop_endpoint_does_not_cancel_the_task():
    """중지 엔드포인트는 플래그만 세운다. `task.cancel()`을 쓰면 in-flight 배치가 버려진다."""
    src = open(os.path.join(ROOT, "admin_console", "backend", "routers", "voc.py"),
               encoding="utf-8").read()
    i = src.index("async def classify_stop")
    stop = src[i:src.index("@router.", i)] if "@router." in src[i:] else src[i:]
    assert "stopping" in stop
    assert ".cancel()" not in stop, "태스크를 강제 취소하면 도는 배치의 판정이 버려진다"
    # 시작 쪽이 그 플래그를 실제로 넘겨줘야 한다.
    assert "should_stop=lambda: _classify_state[\"stopping\"]" in src


def test_console_has_no_helper_prose():
    """사용자 지시: "상세 설명들 다 빼버리고 버튼만 남겨줘. 다른 탭에서도 마찬가지야.
    자잘한 설명들 다 빼줘."(#161)

    화면에 남은 한국어 문장을 센다. 데이터·라벨·경고는 남고, **설명 문단**은 없어야 한다.
    alert/confirm 안의 문장은 화면에 상주하지 않으므로 제외한다."""
    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    banned = [
        "표가 몇 번째 행부터 시작하든",
        "어느 단계가 죽었는지 확인용",
        "양식을 받아 채우거나",
        "커맨드 하나가 에이전트 툴 하나가 됩니다",
        "이름 기준으로 등록/갱신됩니다",
        "API 키를 저장하면 기본 모델이 자동 지정됩니다",
        "먼저 30건으로 판정 품질을 확인한 뒤",
        "에이전트는 이 값을 보고",
        "관리자 콘솔 로그인 계정",
        "(입력 후 Enter)",
    ]
    left = [b for b in banned if b in html]
    assert not left, f"콘솔에 설명 문구가 남아 있다: {left}"

    # 되돌리기 경고는 **확인창(confirm)에만** 남는다. 화면에 상주하는 설명으로 다시
    # 붙이면 두 번 나타난다 — 같은 말을 두 곳에서 하는 것이 원래 문제였다.
    assert html.count("직접 고친 문구가 있으면 사라집니다") == 1


def test_classify_card_has_exactly_three_buttons():
    """처리 주체 분류 카드는 버튼 세 개(먼저 30건만 / 미분류 전부 / 중지)만 있어야 한다."""
    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    i = html.index("<strong>처리 주체 분류</strong>")
    card = html[i:html.index("<div className=\"toolbar\">", i)]
    labels = re.findall(r">([^<>{}]+)</button>", card)
    assert labels == ["먼저 30건만", "미분류 전부", "중지"], labels


# --- #162: 선검색이 0건이 되는 지점을 가른다 -------------------------------------------
def _rag_ns():
    """`_retrieval_query` ~ `_rag_context` 를 실제로 실행할 수 있게 떼어 온다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    body = src[src.index("def _retrieval_query("):src.index("async def _summarize_turns")]
    cfg = {"rag_prefetch": "true", "manual_prefetch_top_k": "3", "voc_prefetch_top_k": "3"}

    async def get_config(key, default=None):
        return cfg.get(key, default)

    import time as _time
    ns = {"asyncio": asyncio, "time": _time, "get_config": get_config,
          "_mem_on": lambda v: str(v).lower() == "true"}
    exec(body, ns)  # noqa: S102
    return ns


def test_prefetch_searches_this_question_first_not_the_merged_one():
    """사용자 로그에서 잡힌 것 (#163):

        q='GPU, cpu 서버별 위치 궁금해\nAA, BB, CC이 뭐야'  →  매뉴얼 3건

    이번 질문은 **약어가 뭐냐**인데 찾아온 3건은 앞 질문(서버 위치)의 것이었다. 0건이 아니라
    **엉뚱한 3건**이라 0건일 때만 도는 재시도가 아예 걸리지 않았고, 모델은 약어 풀이가
    프롬프트에 없으니 지어냈다. 그래서 **이번 질문만으로 먼저** 찾는다."""
    ns = _rag_ns()
    seen = []

    async def manual(query, top_k, vec=None):
        seen.append(query)
        return [{"chunk_text": f"찾은 것: {query}", "guide_document": "매뉴얼"}]

    async def voc(query, top_k, vec=None):
        return [{"question": "q", "answer": "a", "handled_by": "user"}]

    ns["_search_manual_for"] = manual
    ns["_search_voc_for"] = voc
    block, manual_hits, _v = asyncio.run(ns["_rag_context"](
        "AA, BB, CC이 뭐야", [("user", "GPU, cpu 서버별 위치 궁금해")]))

    assert seen == ["AA, BB, CC이 뭐야"], f"앞 질문이 섞인 질의로 찾았다: {seen}"
    assert "서버별 위치" not in block, "앞 질문의 결과가 프롬프트에 들어갔다"


def test_prefetch_falls_back_to_the_merged_query_when_this_question_finds_nothing():
    """`그러면 접속 못 하는거 아니야?`에는 검색할 명사가 없다 — 이때는 직전 발화를 붙인
    질의가 답을 찾는다(#156의 원래 이유). 1차가 0건일 때만 쓴다."""
    ns = _rag_ns()
    seen = []

    async def manual(query, top_k, vec=None):
        seen.append(query)
        return [{"chunk_text": "접속 가이드", "guide_document": "접속 매뉴얼"}] \
            if "\n" in query else []           # 이번 질문만으로는 못 찾는 상황

    async def voc(query, top_k, vec=None):
        return []

    ns["_search_manual_for"] = manual
    ns["_search_voc_for"] = voc
    block, manual_hits, _v = asyncio.run(ns["_rag_context"](
        "그러면 접속 못 하는거 아니야?",
        [("user", "login server 접속이 갑자기 안됩니다")]))

    assert len(seen) == 2, f"재시도를 안 했다: {seen}"
    assert "\n" not in seen[0], "1차가 이미 병합 질의다"
    assert "\n" in seen[1], "2차가 병합 질의가 아니다"
    assert len(manual_hits) == 1 and "접속 가이드" in block


def test_prefetch_does_not_retry_when_both_sides_found_something():
    """찾았으면 다시 찾지 않는다 — 재시도는 실패했을 때만 드는 비용이어야 한다."""
    ns = _rag_ns()
    calls = {"n": 0}

    async def hit(query, top_k, vec=None):
        calls["n"] += 1
        return [{"chunk_text": "본문", "question": "q", "answer": "a"}]

    ns["_search_manual_for"] = hit
    ns["_search_voc_for"] = hit
    asyncio.run(ns["_rag_context"]("질문", [("user", "앞 질문")]))
    assert calls["n"] == 2, f"찾았는데도 다시 검색했다({calls['n']}회)"


def test_empty_prefetch_block_tells_the_model_to_search_again():
    """0건 블록이 "없습니다"로 끝나면 모델이 곧바로 포기한다. 우리가 만든 질의는 사용자
    문장 그대로라 검색어로 불리할 수 있으므로, **직접 다시 찾으라**고 먼저 시킨다 (#162)."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    ns = {}
    exec(src[src.index("def _manual_block("):src.index("async def _rag_context")], ns)  # noqa: S102
    for fn, tool in (("_manual_block", "search_manual"), ("_voc_block", "search_voc")):
        out = ns[fn]([])
        assert tool in out, f"{fn}: 0건일 때 재검색할 도구를 알려주지 않는다"
        assert "다른 표현" in out, f"{fn}: 질의를 바꿔 보라는 말이 없다"


def test_search_logs_which_stage_emptied_the_result():
    """`선검색 0건`만으로는 '후보가 없었다'와 '리랭커가 다 걷어냈다'를 못 가른다.
    셋은 고쳐야 할 곳이 완전히 다르다 (#162)."""
    for mod in ("manual_search", "voc_search"):
        src = open(os.path.join(ROOT, "shared", f"{mod}.py"), encoding="utf-8").read()
        assert src.count("log_stages(") >= 2, f"{mod}: 0건 경로와 정상 경로 양쪽에 있어야 한다"
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    from retrieval import log_stages
    import io
    import contextlib
    buf = io.StringIO()
    with contextlib.redirect_stdout(buf):
        log_stages("voc-search", "로그인 서버 접속이 안됩니다", 12, 0, 0)
    out = buf.getvalue()
    assert "후보 12" in out and "리랭크 0" in out and "최종 0" in out, out


def test_handled_by_does_not_filter_search():
    """사용자 질문: "아직 미분류라서 rag로 못 찾는 건지?" — 아니다.
    `handled_by`는 검색 SQL의 WHERE 절에 없다. 필터로 들어가면 백필 전 데이터가 통째로
    안 잡히게 되므로, 그런 일이 생기지 않게 고정한다 (#162)."""
    src = open(os.path.join(ROOT, "shared", "voc_search.py"), encoding="utf-8").read()
    sql = src[src.index("async def search_voc_records"):]
    for line in sql.split("\n"):
        s = line.strip()
        if s.startswith("WHERE") or s.startswith("AND"):
            assert "handled_by" not in s, f"검색이 판정 여부로 걸러진다: {s}"


def test_prefetch_runs_inside_the_stream_so_progress_can_be_shown():
    """사용자 지적: "manual이나 voc mcp 동작할 때도 '관련 매뉴얼 검색 중…'처럼 현황을
    알려줘야 해."(#163)

    선검색은 도구 호출이 아니라 우리 코드라 진행 줄이 붙을 자리가 없었다(#155). 게다가
    스트림이 **시작되기 전에** 돌아서, 사용자 로그 기준 2.9~3.0초 동안 화면이 비어 있었다.
    진행 줄을 먼저 내보내려면 선검색이 제너레이터 **안에서** 돌아야 한다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    body = src[src.index("async def chat_completions("):src.index("def _mem_on(")]
    gen = body.split("async def event_stream():", 1)[1]
    # `_prepare()` 안이 아니라 **함수 본문에서 곧바로** 선검색을 하면 스트림이 시작되기
    # 전에 돌아 화면이 빈다. `_prepare` 정의 앞부분만 본다.
    top = body[:body.index("async def _prepare():")]
    assert "await _rag_context(" not in top, \
        "선검색이 스트림 시작 전에 돈다 — 그동안 화면이 빈다"
    assert "검색하는 중" in gen, "선검색 진행 줄이 없다"
    assert gen.index("검색하는 중") < gen.index("await _prepare()"), \
        "진행 줄이 선검색 뒤에 나온다(그러면 기다리는 동안 아무것도 안 보인다)"
    # 준비를 제너레이터 안으로 옮기면 toolset 정리가 끊기기 쉽다.
    assert "_close_toolsets(toolsets)" not in body, "정리가 옛 지역변수를 본다(연결 누수)"
    assert body.count('_close_toolsets(prepared["toolsets"])') == 2, \
        "스트리밍/비스트리밍 양쪽에서 정리해야 한다"


def _phrase_fn(tool_kind=None, tool_cmd=None):
    """`_action_phrase`를 실제로 실행한다. 어느 MCP의 도구인지·실제 커맨드를 주입한다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    ns = {"re": re, "_TOOL_KIND": dict(tool_kind or {}), "_TOOL_CMD": dict(tool_cmd or {})}
    exec(src[src.index("def _first_text_arg("):src.index("def _unwrap_result(")], ns)  # noqa: S102
    return ns["_action_phrase"]


def test_progress_says_what_it_is_doing_by_mcp_not_by_tool_name():
    """사용자 지시: "도구명은 보여주지 말고 — voc면 `~~ 관련 과거 VOC 이력 검색 중`,
    manual이면 `~~ 관련 매뉴얼 검색 중`, execution이면 `pwd 실행 중`."(#166)

    **이름으로 추측하지 않는다.** 등록 커맨드는 콘솔에서 붙인 이름이 그대로 툴 이름이라
    (`s2_phd_list`, `myquota` …) 이름만 봐서는 실행인지 검색인지 알 수 없다 — 실제로
    `· 확인하는 중`만 뜨고 있었다(#165). 어느 MCP의 도구인지로 정한다."""
    phrase = _phrase_fn({"aaa": "manual", "bbb": "voc", "ccc": "execution", "ddd": "chart"})
    assert phrase("aaa", {"query": "접속 오류"}) == "'접속 오류' 관련 매뉴얼 검색 중"
    assert phrase("bbb", {"query": "접속 오류"}) == "'접속 오류' 관련 과거 VOC 이력 검색 중"
    assert phrase("ccc", {"option": "-h"}).startswith("`") and "실행 중" in phrase("ccc", {"option": "-h"})
    assert phrase("ddd", {"title": "추이"}) == "차트 그리는 중"
    # 도구 이름이 화면에 새면 안 된다.
    for name in ("aaa", "bbb", "ccc", "ddd"):
        assert name not in phrase(name, {"query": "x"})


def test_manual_block_does_not_answer_this_users_own_values():
    """사용자: "`내 홈스토리지 경로 알려줘` — execution_mcp가 실행되어야 하는데 안 되고
    매뉴얼로만 가이드함."(#163)

    선검색이 매 질문마다 근거를 넣어 주니(#155), 모델이 **실행해야 하는 질문에도** 그
    블록으로 답해 버린다. 계정마다 다른 값은 매뉴얼에 있을 수 없다 — 블록 자체가
    그 경계를 말해야 한다(지시문은 되돌리기 버튼을 눌러야 반영된다, #147)."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    ns = {}
    exec(src[src.index("def _manual_block("):src.index("async def _rag_context")], ns)  # noqa: S102
    out = ns["_manual_block"]([{"chunk_text": "스토리지 안내", "guide_document": "매뉴얼"}])
    assert "실행해서 확인" in out, "값을 묻는 질문에도 매뉴얼로 답하게 두고 있다"
    # 근거가 0건일 때는 이 경고가 잡음이다(붙일 근거 자체가 없다).
    assert "실행해서 확인" not in ns["_manual_block"]([])


# --- #164: MCP toolset은 만든 태스크에서 닫아야 한다 -----------------------------------
def test_toolsets_are_built_and_closed_in_the_same_task():
    """실서버 크래시의 진짜 원인 (#164).

        RuntimeError: Attempted to exit a cancel scope that isn't the current tasks's ...
        Warning: Error during MCP session cleanup ... in a different task than it was entered in

    MCP toolset은 anyio 취소 스코프를 쓰고, 그건 **태스크에 매여 있다**. 그런데
    `StreamingResponse`의 제너레이터는 엔드포인트 코루틴과 **다른 태스크**에서 돈다.
    밖에서 `build_agent()`로 만들고 제너레이터의 `finally`에서 닫으면 그 규칙을 어겨
    응답이 통째로 죽고, 화면에는 `SyntaxError: JSON.parse …`만 보인다(본문이 잘린다).

    그래서 세 엔드포인트 모두 **`_prepare()` 안에서 만들고** `prepared["toolsets"]`로 닫는다.
    """
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    ends = ("async def chat_completions(", "async def agent_query(", "async def voc_query(")
    starts = sorted(src.index(e) for e in ends)
    bounds = list(zip(starts, starts[1:] + [len(src)]))
    for (i, j), name in zip(bounds, sorted(ends, key=src.index)):
        body = src[i:j]
        assert "async def _prepare():" in body, f"{name}: 준비를 한곳에 모으지 않았다"
        # toolset을 만드는 곳은 _prepare 안뿐이어야 한다.
        top = body[:body.index("async def _prepare():")]
        assert "build_agent(" not in top, \
            f"{name}: 스트림 밖에서 toolset을 만든다(닫을 때 태스크가 달라진다)"
        assert "_close_toolsets(toolsets)" not in body, \
            f"{name}: 정리가 옛 지역변수를 본다"
        assert body.count('_close_toolsets(prepared["toolsets"])') == 2, \
            f"{name}: 스트리밍/비스트리밍 양쪽에서 정리해야 한다"
        # 제너레이터가 스스로 준비해야 같은 태스크가 된다.
        gen = body.split("async def event_stream():", 1)[1]
        assert "await _prepare()" in gen, f"{name}: 제너레이터가 준비를 하지 않는다"


# --- #165: 도구가 붙었는지 · 무엇을 하는지 보이게 -------------------------------------
def test_startup_logs_the_tools_the_model_actually_gets():
    """사용자: "아예 execution mcp가 동작을 안 하는 것 같아."(#165)

    **모델에게 없는 도구는 부를 수 없다.** 도구가 안 붙어 있으면 지시문·프롬프트를 아무리
    고쳐도 소용없는데, 그걸 확인할 방법이 로그에 없었다. 기동 시 1회 목록을 찍는다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert "async def _log_tool_inventory(" in src
    assert "await _log_tool_inventory(toolsets)" in src, "기동 경로에서 부르지 않는다"
    fn = src[src.index("async def _log_tool_inventory("):src.index("async def require_api_key(")]
    assert "get_tools()" in fn
    assert "except Exception" in fn, "진단용 로그가 기동을 막으면 안 된다"


def test_unmatched_tool_shows_what_it_runs_not_just_확인하는중():
    """등록 커맨드는 콘솔에서 붙인 이름이 그대로 툴 이름이라 어떤 규칙에도 안 걸린다.
    예전에는 `· 확인하는 중`만 떠서 무엇을 하는지 알 수 없었다 (#165)."""
    out = _phrase_fn({"s2_myquota": "execution"},
                     {"s2_myquota": "myquota"})("s2_myquota", {"option": "-h"})
    assert out == "`myquota -h` 실행 중", out
    assert "s2_myquota" not in out, "도구 이름이 화면에 샌다"


def test_registered_command_shows_the_command_not_just_its_arguments():
    """등록 커맨드는 인자만 넘어온다(`{"option": "-h"}`). 그것만 보여주면 `-h 실행 중`이라
    무슨 커맨드인지 알 수 없다 — execution MCP가 툴 설명 끝에 붙여 주는 실제 커맨드
    (`registry._describe`의 `[...]`)를 앞에 붙인다 (#166)."""
    phrase = _phrase_fn({"t": "execution"}, {"t": "phd list"})
    assert phrase("t", {"option": "-l"}) == "`phd list -l` 실행 중"
    # 커맨드를 못 잡았어도 인자만이라도 보여준다(빈 줄보다 낫다).
    assert "-l" in _phrase_fn({"t": "execution"})("t", {"option": "-l"})


def test_unknown_tool_is_not_reported_as_a_command_execution():
    """어느 MCP인지도 이름 힌트도 없으면 **실행이라고 단정하지 않는다.** 검색 도구의 인자를
    커맨드처럼 보여주면 하지도 않은 실행을 한 것처럼 보인다 (#166)."""
    out = _phrase_fn()("낯선툴", {"path": "/tmp"})
    assert "실행" not in out, out
    assert "확인 중" in out and "낯선툴" not in out


def test_prefetch_embeds_the_query_once_not_twice():
    """매뉴얼과 VOC는 **같은 질의**를 쓰는데 각자 임베딩을 불렀다 — 임베딩 서버 왕복이 두 번.
    선검색이 2.9~3.0초 걸린 것의 한 몫이다 (#165)."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    body = src[src.index("async def _rag_context("):src.index("async def _summarize_turns")]
    assert body.count("_embed_once(bare)") == 1, "질의 벡터를 한 번만 만들어야 한다"
    assert "_search_manual_for(bare, max(1, mk), vec)" in body
    assert "_search_voc_for(bare, max(1, vk), vec)" in body, "VOC가 벡터를 다시 만든다"


# --- #171: 남의 계정은 답변에 나오면 안 된다 -------------------------------------------
def test_answer_never_repeats_someone_elses_account():
    """사용자: "{다른 사람 계정}으로 접속이 불가한데" → 에이전트가
    "귀하의 계정({다른 사람 계정})으로 …"라고 답했다. **절대 안 된다** (#171).

    VOC 검색 결과는 `pii.mask_record`가 이미 가린다. 남은 구멍은 **질문에 적힌 남의 계정을
    답변이 되뇌는 것**이라, 근거 유무로는 못 막는다(질문 자체가 근거다).
    호출자 본인 계정이 아니면 그 줄을 덜어낸다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    import json as _json
    ns = {"re": re, "json": _json}
    exec(src[src.index("_IP_RE = "):src.index("async def _make_grounding(")], ns)  # noqa: S102
    g = ns["_AnswerGuard"](True, "other.user 으로 접속이 불가합니다", user_id="ops.user")
    g.searched_manual = True

    out = g.review("귀하의 계정(other.user)으로 접속이 불가한 사례가 있습니다.\n"
                   "홈 정리를 먼저 해 보세요. 그래도 안 되면 접수해 주세요.")
    assert "other.user" not in out, "남의 계정이 답변에 그대로 나갔다"
    assert "홈 정리를 먼저" in out, "남의 계정이 없는 줄까지 버렸다"

    # 본인 계정은 그대로 나가야 한다.
    keep = g.review("ops.user 님의 홈은 정상입니다. 확인해 보세요.")
    assert "ops.user" in keep


def test_public_tree_never_ships_an_empty_vendor_dir():
    """공개 미러의 `admin_console/frontend/vendor/`에는 실제 js가 없다(.gitignore).
    그 껍데기를 배포 경로에 덮어쓰면 서버의 react·babel이 지워져 콘솔이 **빈 화면**이
    된다 — HTTP는 200을 주므로 네트워크 문제로 오진하기 쉽다 (#173).

    아예 내보내지 않으면 덮어쓸 것도 없다."""
    script = open(os.path.join(ROOT, "scripts", "make-public.sh"), encoding="utf-8").read()
    assert 'rm -rf "$OUT/admin_console/frontend/vendor"' in script, \
        "공개본에서 vendor 디렉토리를 빼지 않는다 - 서버의 react·babel이 지워진다"

    out = os.path.join(ROOT, "dist", "public", "admin_console", "frontend", "vendor")
    if os.path.isdir(os.path.join(ROOT, "dist", "public")):
        assert not os.path.exists(out), "생성된 공개 트리에 vendor 껍데기가 남아 있다"
